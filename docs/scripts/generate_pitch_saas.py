"""
GRANIX SaaS Pitch Deck Generator - XP Investimentos Partnership.

Generates a 15-slide professional PowerPoint presentation for GRANIX's
SaaS model on top of XP Investimentos. Uses dark theme with GRANIX brand
colors and 16:9 aspect ratio.

Usage:
    python generate_pitch_saas.py
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Constants & Brand Identity
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Brand colours
GREEN = RGBColor(0x00, 0xC8, 0x53)
BLUE_DEEP = RGBColor(0x1A, 0x23, 0x7E)
DARK_BG = RGBColor(0x0F, 0x12, 0x20)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xFF, 0x6D, 0x00)
PINK = RGBColor(0xE9, 0x1E, 0x63)
LIGHT_GRAY = RGBColor(0xB0, 0xB0, 0xB0)
MED_GRAY = RGBColor(0x70, 0x70, 0x80)
CARD_BG = RGBColor(0x1A, 0x1E, 0x30)
CARD_BG_ALT = RGBColor(0x22, 0x27, 0x3A)
PURPLE = RGBColor(0x9C, 0x27, 0xB0)
RED = RGBColor(0xF4, 0x43, 0x36)
YELLOW = RGBColor(0xFF, 0xEB, 0x3B)
BLUE = RGBColor(0x21, 0x96, 0xF3)

TITLE_FONT = "Poppins"
BODY_FONT = "Inter"
TITLE_FONT_FALLBACK = "Calibri"
BODY_FONT_FALLBACK = "Arial"

OUTPUT_PATH = "/Users/lab/Claude/Projetos/Granix-App/docs/GRANIX-Pitch-SaaS-XP.pptx"


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------
def set_slide_bg(slide, color: RGBColor) -> None:
    """Set solid background colour on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    font_size: int = 18,
    font_color: RGBColor = WHITE,
    bold: bool = False,
    font_name: str = BODY_FONT,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    """Add a simple text box and return the shape."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = alignment
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    txBox.text_frame._txBody.bodyPr.set(
        "anchor",
        {
            MSO_ANCHOR.TOP: "t",
            MSO_ANCHOR.MIDDLE: "ctr",
            MSO_ANCHOR.BOTTOM: "b",
        }.get(anchor, "t"),
    )
    return txBox


def add_multiline_textbox(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    lines: list[tuple[str, int, RGBColor, bool, str]],
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    line_spacing: float = 1.2,
):
    """Add a textbox with multiple paragraphs (one per line).

    Each line is a tuple: (text, font_size, font_color, bold, font_name).
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, (text, size, color, bld, fname) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bld
        p.font.name = fname
        p.alignment = alignment
        p.space_after = Pt(size * (line_spacing - 1.0))
    return txBox


def add_rich_textbox(
    slide, left: int, top: int, width: int, height: int,
    runs: list[dict], alignment: PP_ALIGN = PP_ALIGN.LEFT,
):
    """Add a text box with multiple styled runs on a single paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    for i, r in enumerate(runs):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
        else:
            run = p.add_run()
        run.text = r.get("text", "")
        run.font.size = Pt(r.get("size", 18))
        run.font.color.rgb = r.get("color", WHITE)
        run.font.bold = r.get("bold", False)
        run.font.name = r.get("font", BODY_FONT)
    return txBox


def add_rounded_rect(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    fill_color: RGBColor,
    border_color: RGBColor | None = None,
    border_width: int = Pt(0),
):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_rect(
    slide, left: int, top: int, width: int, height: int,
    fill_color: RGBColor,
):
    """Add a plain rectangle (no border)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_card(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    title_text: str,
    body_text: str,
    fill: RGBColor = CARD_BG,
    title_color: RGBColor = GREEN,
    body_color: RGBColor = WHITE,
    title_size: int = 16,
    body_size: int = 13,
    border_color: RGBColor | None = None,
    accent_color: RGBColor | None = None,
):
    """Add a card: rounded rect with optional top accent bar, title, and body."""
    add_rounded_rect(
        slide, left, top, width, height, fill,
        border_color=border_color,
        border_width=Pt(2) if border_color else Pt(0),
    )
    if accent_color:
        add_rect(slide, left, top, width, Pt(5), accent_color)
    add_textbox(
        slide, left + Inches(0.2), top + Inches(0.2),
        width - Inches(0.4), Inches(0.4),
        title_text, font_size=title_size, font_color=title_color,
        bold=True, font_name=TITLE_FONT,
    )
    add_textbox(
        slide, left + Inches(0.2), top + Inches(0.65),
        width - Inches(0.4), height - Inches(0.85),
        body_text, font_size=body_size, font_color=body_color,
        font_name=BODY_FONT,
    )


def slide_title(slide, text: str, top: int = Inches(0.3), font_size: int = 32) -> None:
    """Add a prominent slide title."""
    add_textbox(
        slide, Inches(0.8), top, Inches(11.7), Inches(0.8),
        text, font_size=font_size, font_color=WHITE, bold=True,
        font_name=TITLE_FONT,
    )


def slide_number_badge(slide, number: str) -> None:
    """Add a small slide-number badge in the bottom-right corner."""
    add_textbox(
        slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.35),
        number, font_size=10, font_color=MED_GRAY,
        font_name=BODY_FONT, alignment=PP_ALIGN.RIGHT,
    )


def green_accent_line(slide, top: int, left: int = Inches(0.8), width: int = Inches(2.0)) -> None:
    """Add a thin green accent line under a title."""
    add_rect(slide, left, top, width, Pt(3), GREEN)


def stat_box(
    slide, left: int, top: int, width: int, height: int,
    number: str, label: str, num_color: RGBColor = GREEN,
):
    """Add a stat box with a big number and a smaller label below."""
    add_rounded_rect(slide, left, top, width, height, CARD_BG)
    add_textbox(
        slide, left, top + Inches(0.15), width, Inches(0.7),
        number, font_size=38, font_color=num_color,
        bold=True, font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, left + Inches(0.1), top + Inches(0.85),
        width - Inches(0.2), height - Inches(1.0),
        label, font_size=12, font_color=LIGHT_GRAY,
        font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
    )


def check_icon(is_check: bool) -> str:
    """Return a visual indicator for comparison tables."""
    return "\u2713" if is_check else "\u2014"


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def build_slide_01_cover(prs: Presentation) -> None:
    """SLIDE 1 - Cover."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)

    # Top accent bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Pt(4), GREEN)

    # Simulated sprouting plant icon: green circle + stem lines
    cx, cy = Inches(6.666), Inches(1.7)
    r = Inches(0.7)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2
    )
    circle.fill.background()
    circle.line.color.rgb = GREEN
    circle.line.width = Pt(3)

    # Stem
    stem_w, stem_h = Pt(4), Inches(0.35)
    stem = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        cx - stem_w // 2, cy - Inches(0.25),
        stem_w, stem_h,
    )
    stem.fill.solid()
    stem.fill.fore_color.rgb = GREEN
    stem.line.fill.background()

    # Left leaf
    lw, lh = Inches(0.22), Inches(0.10)
    ll = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, cx - Inches(0.22), cy - Inches(0.22), lw, lh
    )
    ll.fill.solid()
    ll.fill.fore_color.rgb = GREEN
    ll.line.fill.background()
    ll.rotation = -40

    # Right leaf
    rl = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, cx + Inches(0.02), cy - Inches(0.22), lw, lh
    )
    rl.fill.solid()
    rl.fill.fore_color.rgb = GREEN
    rl.line.fill.background()
    rl.rotation = 40

    # GRANIX wordmark
    add_rich_textbox(
        slide, Inches(1), Inches(2.6), Inches(11.3), Inches(1.2),
        [
            {"text": "GRAN", "size": 72, "color": WHITE, "bold": True, "font": TITLE_FONT},
            {"text": "IX", "size": 72, "color": GREEN, "bold": True, "font": TITLE_FONT},
        ],
        alignment=PP_ALIGN.CENTER,
    )

    # Tagline
    add_textbox(
        slide, Inches(1.5), Inches(3.8), Inches(10.3), Inches(0.7),
        "Educacao Financeira Gamificada para a Proxima Geracao",
        font_size=24, font_color=LIGHT_GRAY, font_name=BODY_FONT,
        alignment=PP_ALIGN.CENTER,
    )

    # Divider line
    add_rect(slide, Inches(5.5), Inches(4.7), Inches(2.3), Pt(2), GREEN)

    # Partnership subtitle
    add_textbox(
        slide, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.5),
        "Parceria Estrategica com XP Investimentos",
        font_size=20, font_color=GREEN, bold=True,
        font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
    )

    # Confidential
    add_textbox(
        slide, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.4),
        "Confidencial | Marco 2026",
        font_size=12, font_color=MED_GRAY, font_name=BODY_FONT,
        alignment=PP_ALIGN.CENTER,
    )
    slide_number_badge(slide, "1/15")


def build_slide_02_problem(prs: Presentation) -> None:
    """SLIDE 2 - O Problema."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    slide_title(
        slide,
        "85% dos pais querem ensinar educacao financeira. Nenhum app resolve.",
        font_size=28,
    )
    green_accent_line(slide, Inches(1.0))

    # 2x2 stat grid
    stats = [
        ("37.9M", "jovens brasileiros (6-17 anos)\nsem educacao financeira digital", GREEN),
        ("85%", "dos pais interessados em ensinar\neducacao financeira (Serasa 2021)", ORANGE),
        ("93%", "dos jovens 9-17 sao usuarios\nde internet (CETIC 2024)", BLUE),
        ("0", "apps combinam investimento real\n+ educacao gamificada no Brasil", PINK),
    ]

    box_w = Inches(5.5)
    box_h = Inches(1.7)
    gap_x = Inches(0.7)
    gap_y = Inches(0.4)
    x_positions = [Inches(0.8), Inches(0.8) + box_w + gap_x]
    y_positions = [Inches(1.5), Inches(1.5) + box_h + gap_y]

    for i, (num, lbl, color) in enumerate(stats):
        col, row = i % 2, i // 2
        left, top = x_positions[col], y_positions[row]
        add_rounded_rect(slide, left, top, box_w, box_h, CARD_BG)
        # Left accent bar
        add_rect(slide, left, top, Pt(6), box_h, color)
        add_textbox(
            slide, left + Inches(0.3), top + Inches(0.2),
            Inches(2.0), Inches(0.8),
            num, font_size=44, font_color=color, bold=True,
            font_name=TITLE_FONT,
        )
        add_textbox(
            slide, left + Inches(2.4), top + Inches(0.25),
            Inches(2.9), Inches(1.2),
            lbl, font_size=14, font_color=WHITE, font_name=BODY_FONT,
        )

    # Bottom quote
    add_rounded_rect(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.9), CARD_BG_ALT)
    add_textbox(
        slide, Inches(1.2), Inches(5.6), Inches(11.0), Inches(0.7),
        '"A XP tem a melhor conta de investimento para menores. '
        'Mas a experiencia e identica a de um adulto."',
        font_size=15, font_color=LIGHT_GRAY, font_name=BODY_FONT,
        alignment=PP_ALIGN.CENTER,
    )
    slide_number_badge(slide, "2/15")


def build_slide_03_opportunity(prs: Presentation) -> None:
    """SLIDE 3 - A Oportunidade XP."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    slide_title(slide, "XP tem tudo. Menos a experiencia.", font_size=32)
    green_accent_line(slide, Inches(1.0))

    col_w = Inches(5.5)
    col_h = Inches(4.0)
    left_x = Inches(0.8)
    right_x = Inches(6.8)
    top_y = Inches(1.5)

    # LEFT column - O que XP JA TEM
    add_rounded_rect(slide, left_x, top_y, col_w, col_h, CARD_BG, border_color=GREEN, border_width=Pt(2))
    add_textbox(
        slide, left_x + Inches(0.3), top_y + Inches(0.2),
        col_w - Inches(0.6), Inches(0.4),
        "O que XP JA TEM", font_size=18, font_color=GREEN,
        bold=True, font_name=TITLE_FONT,
    )
    xp_has = [
        "Conta completa (PIX, cartao, investimentos)",
        "4.8M clientes ativos",
        "18.200 assessores",
        "R$1.2T sob gestao (AUM)",
        "Banco Modal (BaaS)",
        "Instituto XP",
    ]
    y = top_y + Inches(0.7)
    for item in xp_has:
        add_rich_textbox(
            slide, left_x + Inches(0.3), y, col_w - Inches(0.6), Inches(0.35),
            [
                {"text": "\u2713  ", "size": 14, "color": GREEN, "bold": True, "font": TITLE_FONT},
                {"text": item, "size": 13, "color": WHITE, "font": BODY_FONT},
            ],
        )
        y += Inches(0.42)

    # RIGHT column - O que XP NAO TEM
    add_rounded_rect(slide, right_x, top_y, col_w, col_h, CARD_BG, border_color=ORANGE, border_width=Pt(2))
    add_textbox(
        slide, right_x + Inches(0.3), top_y + Inches(0.2),
        col_w - Inches(0.6), Inches(0.4),
        "O que XP NAO TEM", font_size=18, font_color=ORANGE,
        bold=True, font_name=TITLE_FONT,
    )
    xp_lacks = [
        "App adaptado para jovens",
        "Gamificacao",
        "Educacao financeira interativa",
        "Tarefas / Mesada digital",
        "Dashboard familiar",
        "Engajamento youth",
    ]
    y = top_y + Inches(0.7)
    for item in xp_lacks:
        add_rich_textbox(
            slide, right_x + Inches(0.3), y, col_w - Inches(0.6), Inches(0.35),
            [
                {"text": "\u2716  ", "size": 14, "color": ORANGE, "bold": True, "font": TITLE_FONT},
                {"text": item, "size": 13, "color": WHITE, "font": BODY_FONT},
            ],
        )
        y += Inches(0.42)

    # Bottom warning
    add_rounded_rect(slide, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.7), CARD_BG_ALT)
    add_textbox(
        slide, Inches(1.2), Inches(5.95), Inches(11.0), Inches(0.6),
        "Enquanto isso, Nubank Familia, Inter Kids e C6 Yellow capturam familias.",
        font_size=16, font_color=ORANGE, bold=True,
        font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
    )
    slide_number_badge(slide, "3/15")


def build_slide_04_solution(prs: Presentation) -> None:
    """SLIDE 4 - A Solucao: GRANIX."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    slide_title(slide, "3 Modulos. 1 Plataforma. Sobre a XP.", font_size=32)
    green_accent_line(slide, Inches(1.0))

    modules = [
        (
            "TAREFAS & RECOMPENSAS",
            "Mesada digital, tarefas com recompensa real, "
            "aprovacao dos pais, dinheiro flui para conta XP do menor.",
            ORANGE,
        ),
        (
            "GRANIX ACADEMY",
            "Educacao financeira gamificada, 4 trilhas por idade (6-17), "
            "alinhado BNCC, conectado a investimentos reais na XP.",
            BLUE,
        ),
        (
            "GAMIFICACAO",
            "XP Points, badges, niveis, ranking familiar, "
            "avatar personalizavel, desafios semanais.",
            PURPLE,
        ),
    ]

    box_w = Inches(3.7)
    box_h = Inches(3.8)
    gap = Inches(0.35)
    x_start = Inches(0.8)

    for i, (title, body, accent) in enumerate(modules):
        left = x_start + i * (box_w + gap)
        top = Inches(1.5)
        add_rounded_rect(slide, left, top, box_w, box_h, CARD_BG)
        add_rect(slide, left, top, box_w, Pt(5), accent)

        # Icon placeholder circle
        icon_r = Inches(0.3)
        icon_cx = left + box_w // 2
        icon_cy = top + Inches(0.6)
        icon = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            icon_cx - icon_r, icon_cy - icon_r,
            icon_r * 2, icon_r * 2,
        )
        icon.fill.solid()
        icon.fill.fore_color.rgb = accent
        icon.line.fill.background()

        add_textbox(
            slide, left + Inches(0.2), top + Inches(1.2),
            box_w - Inches(0.4), Inches(0.5),
            title, font_size=16, font_color=accent, bold=True,
            font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, left + Inches(0.25), top + Inches(1.8),
            box_w - Inches(0.5), Inches(1.8),
            body, font_size=13, font_color=WHITE,
            font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
        )

    # Bottom tagline
    add_textbox(
        slide, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.5),
        "Tudo integrado a conta real do menor na XP",
        font_size=16, font_color=GREEN, bold=True,
        font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
    )
    slide_number_badge(slide, "4/15")


def build_slide_05_buckets(prs: Presentation) -> None:
    """SLIDE 5 - 4 Buckets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    slide_title(slide, "O Sistema de 4 Buckets: Unico no Brasil", font_size=30)
    green_accent_line(slide, Inches(1.0))

    buckets = [
        (GREEN, "GASTAR", "40%", "Dinheiro do dia-a-dia\nvia cartao XP"),
        (BLUE, "GUARDAR", "30%", "Objetivos de\ncurto prazo"),
        (PINK, "DOAR", "10%", "Generosidade e\nimpacto social"),
        (PURPLE, "INVESTIR", "20%", "Tesouro Educa+, CDB,\nFundos na XP"),
    ]

    box_w = Inches(2.7)
    box_h = Inches(3.2)
    gap = Inches(0.4)
    x_start = Inches(0.8)

    for i, (color, name, pct, desc) in enumerate(buckets):
        left = x_start + i * (box_w + gap)
        top = Inches(1.4)
        add_rounded_rect(
            slide, left, top, box_w, box_h, CARD_BG,
            border_color=color, border_width=Pt(3),
        )
        add_rect(slide, left, top, box_w, Pt(6), color)

        # Bucket name
        add_textbox(
            slide, left, top + Inches(0.3), box_w, Inches(0.4),
            name, font_size=20, font_color=color, bold=True,
            font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
        )
        # Percentage
        add_textbox(
            slide, left, top + Inches(0.8), box_w, Inches(0.7),
            pct, font_size=42, font_color=WHITE, bold=True,
            font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
        )
        # Description
        add_textbox(
            slide, left + Inches(0.2), top + Inches(1.7),
            box_w - Inches(0.4), Inches(1.2),
            desc, font_size=13, font_color=LIGHT_GRAY,
            font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
        )

    # Bottom text
    add_textbox(
        slide, Inches(0.8), Inches(4.9), Inches(11.7), Inches(0.5),
        "Cada real que entra e distribuido automaticamente nos 4 buckets",
        font_size=15, font_color=WHITE, font_name=BODY_FONT,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.4),
        "Nenhum concorrente brasileiro oferece este sistema",
        font_size=15, font_color=GREEN, bold=True,
        font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
    )
    slide_number_badge(slide, "5/15")


def build_slide_06_greenlight(prs: Presentation) -> None:
    """SLIDE 6 - Benchmark: Greenlight."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    slide_title(slide, "Greenlight provou o modelo nos EUA", font_size=32)
    green_accent_line(slide, Inches(1.0))

    metrics = [
        ("US$228.5M", "revenue (2024)"),
        ("$3.2B", "valuation"),
        ("6.5M", "familias ativas"),
        ("150+", "bancos parceiros"),
        ("US$556M", "funding total"),
        ("JPMorgan", "Wells Fargo,\nMorgan Stanley"),
    ]

    cols = 3
    box_w = Inches(3.6)
    box_h = Inches(1.5)
    gap_x = Inches(0.35)
    gap_y = Inches(0.3)
    x_start = Inches(0.8)
    y_start = Inches(1.5)

    for i, (num, lbl) in enumerate(metrics):
        r = i // cols
        c = i % cols
        stat_box(
            slide,
            x_start + c * (box_w + gap_x),
            y_start + r * (box_h + gap_y),
            box_w, box_h, num, lbl,
        )

    # Key insight
    add_rounded_rect(slide, Inches(0.8), Inches(4.9), Inches(11.7), Inches(0.8), CARD_BG_ALT)
    add_textbox(
        slide, Inches(1.2), Inches(4.95), Inches(11.0), Inches(0.7),
        "Greenlight for Banks: bancos PAGAM pelo acesso a plataforma",
        font_size=18, font_color=GREEN, bold=True,
        font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.5),
        "GRANIX e o Greenlight brasileiro. Sobre a XP.",
        font_size=20, font_color=WHITE, bold=True,
        font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
    )
    slide_number_badge(slide, "6/15")


def build_slide_07_market(prs: Presentation) -> None:
    """SLIDE 7 - Mercado Brasileiro."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    slide_title(slide, "Mercado de R$8.1B com zero solucao completa", font_size=30)
    green_accent_line(slide, Inches(1.0))

    # Concentric circles (using progressively smaller rounded rects)
    circles_data = [
        ("TAM", "R$8.1B", "25.9M familias com filhos\nR$216/ano", Inches(5.0), Inches(3.6), BLUE_DEEP),
        ("SAM", "R$1.45B", "6.7M familias A/B/C\ncom smartphone", Inches(3.6), Inches(2.8), RGBColor(0x25, 0x2E, 0x80)),
        ("SOM", "500K fam.", "em 3 anos\n7.5% do SAM", Inches(2.2), Inches(2.0), RGBColor(0x00, 0x5A, 0x25)),
    ]

    center_x = Inches(3.8)
    center_y = Inches(3.8)

    for label, value, desc, w, h, color in circles_data:
        left = center_x - w // 2
        top = center_y - h // 2
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, left, top, w, h,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = GREEN
        shape.line.width = Pt(2)

    # Labels on the right side
    label_x = Inches(7.0)
    labels = [
        ("TAM: R$8.1B", "25.9M familias com filhos, R$216/ano", Inches(1.8), GREEN),
        ("SAM: R$1.45B", "6.7M familias classes A/B/C com smartphone", Inches(3.0), BLUE),
        ("SOM: 500K familias", "em 3 anos (7.5% do SAM)", Inches(4.2), ORANGE),
    ]
    for title, desc, y, color in labels:
        add_rounded_rect(slide, label_x, y, Inches(5.5), Inches(0.9), CARD_BG)
        add_rect(slide, label_x, y, Pt(6), Inches(0.9), color)
        add_textbox(
            slide, label_x + Inches(0.3), y + Inches(0.08),
            Inches(5.0), Inches(0.35),
            title, font_size=16, font_color=color, bold=True,
            font_name=TITLE_FONT,
        )
        add_textbox(
            slide, label_x + Inches(0.3), y + Inches(0.45),
            Inches(5.0), Inches(0.35),
            desc, font_size=12, font_color=LIGHT_GRAY, font_name=BODY_FONT,
        )

    slide_number_badge(slide, "7/15")


def build_slide_08_competitors(prs: Presentation) -> None:
    """SLIDE 8 - Competidores."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    slide_title(slide, "Ninguem combina investimento + educacao gamificada", font_size=28)
    green_accent_line(slide, Inches(1.0))

    headers = ["", "Cartao", "Tarefas", "Educacao", "Investimento", "Gamificacao", "4 Buckets"]
    # Each row: name, then True/False for each feature
    rows_data = [
        ("Nubank Familia",  True,  False, False, False, False, False),
        ("Inter Kids",      True,  False, False, True,  False, False),
        ("C6 Yellow",       True,  False, False, True,  False, False),
        ("NextJoy",         True,  True,  False, False, False, False),
        ("BTG Jovem",       True,  False, False, True,  False, False),
        ("XP (sozinha)",    True,  False, False, True,  False, False),
        ("XP + GRANIX",     True,  True,  True,  True,  True,  True),
    ]

    col_widths = [Inches(2.2), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.7), Inches(1.7), Inches(1.6)]
    x_start = Inches(0.4)
    y_start = Inches(1.4)
    row_h = Inches(0.55)

    # Header row
    x = x_start
    for j, h in enumerate(headers):
        add_rounded_rect(slide, x, y_start, col_widths[j] - Pt(3), row_h - Pt(3), BLUE_DEEP)
        add_textbox(
            slide, x + Pt(4), y_start + Pt(2),
            col_widths[j] - Pt(12), row_h - Pt(8),
            h, font_size=11, font_color=GREEN, bold=True,
            font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        x += col_widths[j]

    # Data rows
    for i, row in enumerate(rows_data):
        name = row[0]
        features = row[1:]
        is_granix = (i == len(rows_data) - 1)
        y = y_start + (i + 1) * row_h
        x = x_start

        for j in range(len(headers)):
            if is_granix:
                bg = RGBColor(0x00, 0x3D, 0x1A)
            elif i % 2 == 0:
                bg = CARD_BG
            else:
                bg = CARD_BG_ALT

            add_rounded_rect(slide, x, y, col_widths[j] - Pt(3), row_h - Pt(3), bg)

            if j == 0:
                cell_text = name
                cell_color = GREEN if is_granix else WHITE
                cell_bold = is_granix
            else:
                has_feature = features[j - 1]
                cell_text = check_icon(has_feature)
                if has_feature:
                    cell_color = GREEN
                else:
                    cell_color = RGBColor(0x55, 0x55, 0x55)
                cell_bold = is_granix and has_feature

            add_textbox(
                slide, x + Pt(4), y + Pt(2),
                col_widths[j] - Pt(12), row_h - Pt(8),
                cell_text, font_size=12 if j > 0 else 11,
                font_color=cell_color, bold=cell_bold,
                font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            x += col_widths[j]

    # Bottom
    add_textbox(
        slide, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.5),
        "XP + GRANIX = unica solucao completa no mercado",
        font_size=16, font_color=GREEN, bold=True,
        font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
    )
    slide_number_badge(slide, "8/15")


def build_slide_09_business_model(prs: Presentation) -> None:
    """SLIDE 9 - Modelo de Negocio."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    slide_title(slide, "Subscription Only. Simples e Escalavel.", font_size=30)
    green_accent_line(slide, Inches(1.0))

    plans = [
        ("GRATIS", "R$0", "/mes",
         "Tarefas basicas\nAcademy limitada\n2 buckets",
         CARD_BG, None),
        ("FAMILIA", "R$19,90", "/mes",
         "Tudo ilimitado\n4 buckets completos\nGamificacao completa\nDashboard familiar",
         CARD_BG, GREEN),
        ("FAMILIA+", "R$39,90", "/mes",
         "Multiplos filhos\nRelatorios avancados\nSimuladores de investimento\nConteudo exclusivo",
         CARD_BG, None),
    ]

    box_w = Inches(3.6)
    box_h = Inches(3.8)
    gap = Inches(0.4)
    x_start = Inches(0.8)

    for i, (name, price, period, features, bg, highlight) in enumerate(plans):
        left = x_start + i * (box_w + gap)
        top = Inches(1.4)
        border = highlight
        add_rounded_rect(
            slide, left, top, box_w, box_h, bg,
            border_color=border, border_width=Pt(3) if border else Pt(0),
        )

        if highlight:
            tag = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left + box_w - Inches(1.6), top - Inches(0.2),
                Inches(1.6), Inches(0.35),
            )
            tag.fill.solid()
            tag.fill.fore_color.rgb = GREEN
            tag.line.fill.background()
            add_textbox(
                slide, left + box_w - Inches(1.6), top - Inches(0.18),
                Inches(1.6), Inches(0.3),
                "RECOMENDADO", font_size=10, font_color=WHITE,
                bold=True, font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
            )

        add_textbox(
            slide, left, top + Inches(0.25), box_w, Inches(0.4),
            name, font_size=18, font_color=GREEN if not highlight else WHITE,
            bold=True, font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
        )

        add_rich_textbox(
            slide, left, top + Inches(0.75), box_w, Inches(0.6),
            [
                {"text": price, "size": 32, "color": WHITE, "bold": True, "font": TITLE_FONT},
                {"text": period, "size": 14, "color": LIGHT_GRAY, "font": BODY_FONT},
            ],
            alignment=PP_ALIGN.CENTER,
        )

        add_textbox(
            slide, left + Inches(0.3), top + Inches(1.6),
            box_w - Inches(0.6), Inches(2.0),
            features, font_size=13, font_color=LIGHT_GRAY,
            font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
        )

    # Revenue share visual
    share_y = Inches(5.6)
    add_rounded_rect(slide, Inches(0.8), share_y, Inches(11.7), Inches(1.0), CARD_BG_ALT)

    # 70% GRANIX bar
    bar_left = Inches(1.5)
    bar_top = share_y + Inches(0.2)
    bar_h = Inches(0.55)
    granix_w = Inches(7.0)
    xp_w = Inches(3.0)

    add_rounded_rect(slide, bar_left, bar_top, granix_w, bar_h, GREEN)
    add_textbox(
        slide, bar_left, bar_top, granix_w, bar_h,
        "70% GRANIX", font_size=14, font_color=WHITE,
        bold=True, font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    add_rounded_rect(slide, bar_left + granix_w + Pt(3), bar_top, xp_w, bar_h, BLUE_DEEP)
    add_textbox(
        slide, bar_left + granix_w + Pt(3), bar_top, xp_w, bar_h,
        "30% XP", font_size=14, font_color=WHITE,
        bold=True, font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    add_textbox(
        slide, Inches(0.8), share_y + Inches(0.75), Inches(11.7), Inches(0.3),
        "XP tambem retem 100% do interchange, float e spread de investimentos",
        font_size=11, font_color=MED_GRAY, font_name=BODY_FONT,
        alignment=PP_ALIGN.CENTER,
    )
    slide_number_badge(slide, "9/15")


def build_slide_10_why_xp(prs: Presentation) -> None:
    """SLIDE 10 - Por que XP diz SIM."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    slide_title(slide, "5 Razoes para a XP", font_size=32)
    green_accent_line(slide, Inches(1.0))

    reasons = [
        (
            "1. CAC Intergeracional",
            "R$0 vs R$500-1,000/adulto. Crianca que cresce na XP vira cliente fiel aos 18.",
            GREEN,
        ),
        (
            "2. Defesa Competitiva",
            "Nubank, Inter, C6 ja tem produtos youth. XP precisa responder agora.",
            ORANGE,
        ),
        (
            "3. Instituto XP",
            "Produtiza a missao de educacao financeira do Instituto com tecnologia.",
            BLUE,
        ),
        (
            "4. Revenue Share",
            "30% da receita recorrente sem nenhum custo de desenvolvimento.",
            PINK,
        ),
        (
            "5. Greenlight Provou",
            "150+ bancos parceiros no modelo identico. JPMorgan, Wells Fargo, Morgan Stanley.",
            PURPLE,
        ),
    ]

    # Layout: 3 cards on top row, 2 on bottom row (centered)
    box_w = Inches(3.6)
    box_h = Inches(2.0)
    gap = Inches(0.4)

    for i, (title, body, accent) in enumerate(reasons):
        if i < 3:
            left = Inches(0.8) + i * (box_w + gap)
            top = Inches(1.4)
        else:
            left = Inches(0.8) + (box_w + gap) * 0.5 + (i - 3) * (box_w + gap)
            top = Inches(3.7)

        add_card(
            slide, left, top, box_w, box_h,
            title, body,
            title_color=accent, body_size=12,
            accent_color=accent,
        )

    slide_number_badge(slide, "10/15")


def build_slide_11_cac(prs: Presentation) -> None:
    """SLIDE 11 - CAC Intergeracional (deep dive)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    slide_title(slide, "O cliente que custa R$0", font_size=32)
    green_accent_line(slide, Inches(1.0))

    col_w = Inches(5.5)
    col_h = Inches(3.2)
    left_x = Inches(0.8)
    right_x = Inches(6.8)
    top_y = Inches(1.5)

    # LEFT - Caminho Tradicional
    add_rounded_rect(slide, left_x, top_y, col_w, col_h, CARD_BG, border_color=RED, border_width=Pt(2))
    add_textbox(
        slide, left_x + Inches(0.3), top_y + Inches(0.2),
        col_w - Inches(0.6), Inches(0.4),
        "Caminho Tradicional", font_size=18, font_color=RED,
        bold=True, font_name=TITLE_FONT,
    )
    trad_lines = [
        "Marketing R$500-1,000",
        "\u2193",
        "1 cliente adulto",
        "\u2193",
        "Pode churnar a qualquer momento",
    ]
    y = top_y + Inches(0.7)
    for line in trad_lines:
        color = RED if line == "\u2193" else WHITE
        add_textbox(
            slide, left_x + Inches(0.3), y, col_w - Inches(0.6), Inches(0.35),
            line, font_size=14, font_color=color,
            font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
        )
        y += Inches(0.38)

    # RIGHT - Caminho GRANIX
    add_rounded_rect(slide, right_x, top_y, col_w, col_h, CARD_BG, border_color=GREEN, border_width=Pt(2))
    add_textbox(
        slide, right_x + Inches(0.3), top_y + Inches(0.2),
        col_w - Inches(0.6), Inches(0.4),
        "Caminho GRANIX", font_size=18, font_color=GREEN,
        bold=True, font_name=TITLE_FONT,
    )
    granix_lines = [
        "Crianca 10 anos entra via pai XP",
        "\u2193",
        "8 anos de educacao financeira",
        "\u2193",
        "Cliente fiel aos 18 | CAC = R$0",
    ]
    y = top_y + Inches(0.7)
    for line in granix_lines:
        color = GREEN if line == "\u2193" else WHITE
        add_textbox(
            slide, right_x + Inches(0.3), y, col_w - Inches(0.6), Inches(0.35),
            line, font_size=14, font_color=color,
            font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
        )
        y += Inches(0.38)

    # Math callout
    add_rounded_rect(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.7), CARD_BG_ALT)
    add_textbox(
        slide, Inches(1.2), Inches(5.05), Inches(11.0), Inches(0.6),
        "Com 4.8M clientes, 5% de adocao = 240.000 familias = pipeline de futuros investidores",
        font_size=16, font_color=WHITE, bold=True,
        font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.5),
        "Greenlight provou: JPMorgan e Wells Fargo investiram E se tornaram parceiros",
        font_size=14, font_color=LIGHT_GRAY,
        font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
    )
    slide_number_badge(slide, "11/15")


def build_slide_12_financials(prs: Presentation) -> None:
    """SLIDE 12 - Projecoes Financeiras."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    slide_title(slide, "Projecao 5 Anos - Cenario com Base XP", font_size=28)
    green_accent_line(slide, Inches(1.0))

    # Table data
    headers = ["", "Ano 1", "Ano 2", "Ano 3", "Ano 4", "Ano 5"]
    rows_data = [
        ("Familias totais",      "48K",     "144K",    "240K",    "336K",    "480K"),
        ("Familias pagantes",    "24K",     "72K",     "120K",    "168K",    "240K"),
        ("Revenue bruto",        "R$6.6M",  "R$21.5M", "R$38.7M", "R$58.2M", "R$86.1M"),
        ("Share Granix (70%)",   "R$4.6M",  "R$15.1M", "R$27.1M", "R$40.7M", "R$60.3M"),
        ("Share XP (30%)",       "R$2.0M",  "R$6.5M",  "R$11.6M", "R$17.5M", "R$25.8M"),
    ]

    col_widths = [Inches(2.8)] + [Inches(1.9)] * 5
    x_start = Inches(0.5)
    y_start = Inches(1.4)
    row_h = Inches(0.52)

    # Header row
    x = x_start
    for j, h in enumerate(headers):
        add_rounded_rect(slide, x, y_start, col_widths[j] - Pt(3), row_h - Pt(3), BLUE_DEEP)
        add_textbox(
            slide, x + Pt(4), y_start + Pt(2),
            col_widths[j] - Pt(12), row_h - Pt(8),
            h, font_size=12, font_color=GREEN, bold=True,
            font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        x += col_widths[j]

    # Data rows
    for i, (label, *values) in enumerate(rows_data):
        y = y_start + (i + 1) * row_h
        x = x_start
        all_cells = [label] + list(values)
        for j, cell in enumerate(all_cells):
            bg = CARD_BG if i % 2 == 0 else CARD_BG_ALT
            # Highlight Granix share row
            if i == 3:
                bg = RGBColor(0x00, 0x3D, 0x1A)
            add_rounded_rect(slide, x, y, col_widths[j] - Pt(3), row_h - Pt(3), bg)
            cell_color = LIGHT_GRAY if j == 0 else WHITE
            if i == 3 and j > 0:
                cell_color = GREEN
            add_textbox(
                slide, x + Pt(4), y + Pt(2),
                col_widths[j] - Pt(12), row_h - Pt(8),
                cell, font_size=12, font_color=cell_color,
                bold=(i == 3), font_name=BODY_FONT,
                alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            x += col_widths[j]

    # Simple bar chart visualization
    chart_y = Inches(4.5)
    chart_h = Inches(2.0)
    bar_area_h = Inches(1.5)
    bar_y_base = chart_y + chart_h - Inches(0.2)

    add_textbox(
        slide, Inches(0.8), chart_y, Inches(3.0), Inches(0.4),
        "Revenue Bruto (R$M)", font_size=12, font_color=LIGHT_GRAY,
        bold=True, font_name=TITLE_FONT,
    )

    revenues = [6.6, 21.5, 38.7, 58.2, 86.1]
    max_rev = 86.1
    bar_w = Inches(1.5)
    bar_gap = Inches(0.5)
    bar_x_start = Inches(1.5)

    for i, rev in enumerate(revenues):
        bh = int(bar_area_h * (rev / max_rev))
        bx = bar_x_start + i * (bar_w + bar_gap)
        by = bar_y_base - bh

        add_rounded_rect(slide, bx, by, bar_w, bh, GREEN)
        # Value on top
        add_textbox(
            slide, bx, by - Inches(0.3), bar_w, Inches(0.3),
            f"R${rev}M", font_size=10, font_color=GREEN,
            bold=True, font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
        )
        # Year label below
        add_textbox(
            slide, bx, bar_y_base + Pt(2), bar_w, Inches(0.25),
            f"Ano {i + 1}", font_size=10, font_color=LIGHT_GRAY,
            font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
        )

    slide_number_badge(slide, "12/15")


def build_slide_13_unit_economics(prs: Presentation) -> None:
    """SLIDE 13 - Unit Economics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    slide_title(slide, "Unit Economics Saudaveis", font_size=32)
    green_accent_line(slide, Inches(1.0))

    metrics = [
        ("R$22,90", "ARPU/mes\n(pagantes)", GREEN),
        ("R$85", "CAC\n(via assessores XP)", ORANGE),
        ("R$687", "LTV\n(30 meses avg)", BLUE),
        ("8.1x", "LTV / CAC", GREEN),
        ("3.7", "meses\nPayback", PURPLE),
    ]

    box_w = Inches(2.2)
    box_h = Inches(2.5)
    gap = Inches(0.25)
    x_start = Inches(0.6)

    for i, (num, lbl, color) in enumerate(metrics):
        left = x_start + i * (box_w + gap)
        top = Inches(1.5)
        add_rounded_rect(slide, left, top, box_w, box_h, CARD_BG, border_color=color, border_width=Pt(2))
        add_textbox(
            slide, left, top + Inches(0.3), box_w, Inches(0.8),
            num, font_size=36, font_color=color, bold=True,
            font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, left + Inches(0.15), top + Inches(1.3),
            box_w - Inches(0.3), Inches(1.0),
            lbl, font_size=13, font_color=LIGHT_GRAY,
            font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
        )

    # Benchmark comparison
    add_rounded_rect(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(1.0), CARD_BG_ALT)
    add_textbox(
        slide, Inches(1.2), Inches(4.55), Inches(11.0), Inches(0.9),
        "Benchmark: Greenlight LTV/CAC ~5x. GRANIX supera por ter CAC mais baixo\n"
        "via distribuicao XP (assessores + base de clientes existente).",
        font_size=14, font_color=LIGHT_GRAY,
        font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
    )
    slide_number_badge(slide, "13/15")


def build_slide_14_roadmap(prs: Presentation) -> None:
    """SLIDE 14 - Roadmap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    slide_title(slide, "18 Meses ate Escala", font_size=32)
    green_accent_line(slide, Inches(1.0))

    phases = [
        (
            "FASE 1", "M1-4", "MVP",
            "Tarefas + 10 licoes\nGamificacao basica\nIntegracao conta XP",
            GREEN,
        ),
        (
            "FASE 2", "M5-8", "Academy Completa",
            "Academy 4 trilhas\nSimuladores\nTesouro Educa+ integration",
            BLUE,
        ),
        (
            "FASE 3", "M9-12", "Gamificacao Avancada",
            "Avatar personalizavel\nRanking + desafios\nCartao Modal",
            ORANGE,
        ),
        (
            "FASE 4", "M13-18", "Escala",
            "White-label\nExpansao LATAM\nParcerias escolas",
            PURPLE,
        ),
    ]

    box_w = Inches(2.8)
    box_h = Inches(3.5)
    gap = Inches(0.3)
    x_start = Inches(0.5)

    for i, (phase, months, subtitle, desc, color) in enumerate(phases):
        left = x_start + i * (box_w + gap)
        top = Inches(1.5)

        add_rounded_rect(slide, left, top, box_w, box_h, CARD_BG)
        add_rect(slide, left, top, box_w, Pt(5), color)

        # Phase badge
        badge_w = Inches(1.2)
        badge_h = Inches(0.35)
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left + (box_w - badge_w) // 2, top + Inches(0.25),
            badge_w, badge_h,
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        add_textbox(
            slide, left + (box_w - badge_w) // 2, top + Inches(0.25),
            badge_w, badge_h,
            phase, font_size=12, font_color=WHITE,
            bold=True, font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        # Months
        add_textbox(
            slide, left, top + Inches(0.75), box_w, Inches(0.35),
            months, font_size=16, font_color=color,
            bold=True, font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
        )

        # Subtitle
        add_textbox(
            slide, left, top + Inches(1.15), box_w, Inches(0.35),
            subtitle, font_size=14, font_color=WHITE,
            bold=True, font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
        )

        # Description
        add_textbox(
            slide, left + Inches(0.2), top + Inches(1.7),
            box_w - Inches(0.4), Inches(1.6),
            desc, font_size=12, font_color=LIGHT_GRAY,
            font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
        )

    # Timeline connector line
    line_y = Inches(5.3)
    add_rect(
        slide, Inches(0.5), line_y,
        Inches(12.3), Pt(3), GREEN,
    )
    # Dots on the line
    for i in range(4):
        dot_x = Inches(0.5) + Inches(0.5) + i * (box_w + gap) + box_w // 2 - Inches(0.08)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, dot_x, line_y - Inches(0.06),
            Inches(0.16), Inches(0.16),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = GREEN
        dot.line.fill.background()

    slide_number_badge(slide, "14/15")


def build_slide_15_ask(prs: Presentation) -> None:
    """SLIDE 15 - O Ask."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    slide_title(slide, "Proposta de Parceria", font_size=32)
    green_accent_line(slide, Inches(1.0))

    asks = [
        (
            "1. Acesso a APIs",
            "APIs da XP e Banco Modal\npara integracao completa\n(contas, PIX, investimentos)",
            BLUE,
        ),
        (
            "2. Distribuicao",
            "Via assessores e base\nde clientes XP\n(4.8M familias potenciais)",
            GREEN,
        ),
        (
            "3. Pre-Seed R$1.5M",
            "Equity a negociar\nRunway de 18 meses\nMVP + lancamento + escala",
            ORANGE,
        ),
    ]

    box_w = Inches(3.6)
    box_h = Inches(2.2)
    gap = Inches(0.4)
    x_start = Inches(0.8)

    for i, (title, body, accent) in enumerate(asks):
        left = x_start + i * (box_w + gap)
        top = Inches(1.4)
        add_rounded_rect(
            slide, left, top, box_w, box_h, CARD_BG,
            border_color=accent, border_width=Pt(2),
        )
        add_rect(slide, left, top, box_w, Pt(5), accent)
        add_textbox(
            slide, left + Inches(0.25), top + Inches(0.25),
            box_w - Inches(0.5), Inches(0.5),
            title, font_size=16, font_color=accent, bold=True,
            font_name=TITLE_FONT,
        )
        add_textbox(
            slide, left + Inches(0.25), top + Inches(0.8),
            box_w - Inches(0.5), Inches(1.2),
            body, font_size=13, font_color=WHITE,
            font_name=BODY_FONT,
        )

    # Revenue share reminder
    share_y = Inches(3.9)
    add_rounded_rect(slide, Inches(0.8), share_y, Inches(11.7), Inches(0.6), CARD_BG_ALT)
    add_textbox(
        slide, Inches(1.2), share_y + Inches(0.1), Inches(11.0), Inches(0.4),
        "Revenue Share: 70% GRANIX / 30% XP",
        font_size=16, font_color=GREEN, bold=True,
        font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
    )

    # Big closing quote
    add_rounded_rect(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.0), CARD_BG)
    add_textbox(
        slide, Inches(1.2), Inches(4.85), Inches(11.0), Inches(0.9),
        '"A questao nao e SE alguem vai construir o Greenlight brasileiro.\n'
        'E QUEM e QUANDO."',
        font_size=22, font_color=WHITE, bold=True,
        font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
    )

    # Contact
    add_textbox(
        slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.4),
        "parceiros@granix.app",
        font_size=16, font_color=GREEN,
        font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
    )

    # Obrigado
    add_textbox(
        slide, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.4),
        "Obrigado.",
        font_size=18, font_color=WHITE, bold=True,
        font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
    )

    # Confidential
    add_textbox(
        slide, Inches(9.5), Inches(7.0), Inches(3.5), Inches(0.35),
        "Confidencial | Marco 2026", font_size=10, font_color=MED_GRAY,
        font_name=BODY_FONT, alignment=PP_ALIGN.RIGHT,
    )
    slide_number_badge(slide, "15/15")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main() -> None:
    """Generate the GRANIX SaaS pitch deck."""
    prs = Presentation()

    # Set 16:9 aspect ratio
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Build all 15 slides
    build_slide_01_cover(prs)
    build_slide_02_problem(prs)
    build_slide_03_opportunity(prs)
    build_slide_04_solution(prs)
    build_slide_05_buckets(prs)
    build_slide_06_greenlight(prs)
    build_slide_07_market(prs)
    build_slide_08_competitors(prs)
    build_slide_09_business_model(prs)
    build_slide_10_why_xp(prs)
    build_slide_11_cac(prs)
    build_slide_12_financials(prs)
    build_slide_13_unit_economics(prs)
    build_slide_14_roadmap(prs)
    build_slide_15_ask(prs)

    prs.save(OUTPUT_PATH)
    print(f"Generated {len(prs.slides)} slides -> {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
