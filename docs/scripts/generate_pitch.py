"""
GRANIX Pitch Deck Generator for XP/Modal Partnership.

Generates a 15-slide professional PowerPoint presentation using python-pptx.
Uses dark theme with GRANIX brand colors.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Brand colours
GREEN = RGBColor(0x00, 0xC8, 0x53)
BLUE_DEEP = RGBColor(0x1A, 0x23, 0x7E)
ORANGE = RGBColor(0xFF, 0x6D, 0x00)
BLUE = RGBColor(0x21, 0x96, 0xF3)
PINK = RGBColor(0xE9, 0x1E, 0x63)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB0, 0xB0)
DARK_BG = RGBColor(0x1E, 0x1E, 0x2E)
DARK_BG2 = RGBColor(0x2A, 0x2A, 0x3E)
CARD_BG = RGBColor(0x2D, 0x2D, 0x44)
RED = RGBColor(0xF4, 0x43, 0x36)
YELLOW = RGBColor(0xFF, 0xEB, 0x3B)
CARD_HIGHLIGHT = RGBColor(0x00, 0x3D, 0x1A)

TITLE_FONT = "Poppins"
BODY_FONT = "Inter"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def set_slide_bg(slide, color: RGBColor):
    """Set solid background colour on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    font_color=WHITE,
    bold=False,
    font_name=BODY_FONT,
    alignment=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    """Add a simple text box and return the shape."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = alignment
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    # vertical anchor
    txBox.text_frame._txBody.bodyPr.set("anchor", {
        MSO_ANCHOR.TOP: "t",
        MSO_ANCHOR.MIDDLE: "ctr",
        MSO_ANCHOR.BOTTOM: "b",
    }.get(anchor, "t"))
    return txBox


def add_rich_textbox(slide, left, top, width, height, runs, alignment=PP_ALIGN.LEFT):
    """Add a text box with multiple styled runs on a single paragraph.

    runs: list of dicts with keys: text, size, color, bold, font
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    for i, r in enumerate(runs):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
        else:
            run = p.add_run()
        run.text = r.get("text", "")
        run.font.size = Pt(r.get("size", 18))
        run.font.color.rgb = r.get("color", WHITE)
        run.font.bold = r.get("bold", False)
        run.font.name = r.get("font", BODY_FONT)
    return txBox


def add_rounded_rect(
    slide,
    left,
    top,
    width,
    height,
    fill_color,
    border_color=None,
    border_width=Pt(0),
):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_card(
    slide,
    left,
    top,
    width,
    height,
    title_text,
    body_text,
    fill=CARD_BG,
    title_color=GREEN,
    body_color=WHITE,
    title_size=16,
    body_size=13,
    border_color=None,
):
    """Add a card (rounded rect with title and body text)."""
    add_rounded_rect(slide, left, top, width, height, fill, border_color=border_color)
    add_textbox(
        slide, left + Inches(0.2), top + Inches(0.15),
        width - Inches(0.4), Inches(0.4),
        title_text, font_size=title_size, font_color=title_color,
        bold=True, font_name=TITLE_FONT, alignment=PP_ALIGN.LEFT,
    )
    add_textbox(
        slide, left + Inches(0.2), top + Inches(0.55),
        width - Inches(0.4), height - Inches(0.7),
        body_text, font_size=body_size, font_color=body_color,
        font_name=BODY_FONT, alignment=PP_ALIGN.LEFT,
    )


def slide_title(slide, text, top=Inches(0.3), font_size=36):
    """Add a slide title."""
    add_textbox(
        slide, Inches(0.8), top, Inches(11.7), Inches(0.7),
        text, font_size=font_size, font_color=WHITE, bold=True,
        font_name=TITLE_FONT, alignment=PP_ALIGN.LEFT,
    )


def slide_subtitle(slide, text, top=Inches(0.9), font_size=18, color=LIGHT_GRAY):
    """Add a slide subtitle."""
    add_textbox(
        slide, Inches(0.8), top, Inches(11.7), Inches(0.5),
        text, font_size=font_size, font_color=color,
        font_name=BODY_FONT, alignment=PP_ALIGN.LEFT,
    )


def stat_box(slide, left, top, width, height, number, label, num_color=GREEN):
    """Add a stat box with a big number and label."""
    add_rounded_rect(slide, left, top, width, height, CARD_BG)
    add_textbox(
        slide, left, top + Inches(0.15), width, Inches(0.7),
        number, font_size=42, font_color=num_color,
        bold=True, font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, left + Inches(0.15), top + Inches(0.85), width - Inches(0.3), height - Inches(1.0),
        label, font_size=12, font_color=LIGHT_GRAY,
        font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
    )


def add_green_line(slide, top):
    """Add a thin green accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), top, Inches(1.5), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def draw_logo_icon(slide, cx_inches, cy_inches, size_inches):
    """Draw the GRANIX sprout logo icon using shapes (circle + seed + stem + leaves)."""
    from pptx.util import Emu
    s = Inches(size_inches)
    left = Inches(cx_inches) - s // 2
    top = Inches(cy_inches) - s // 2

    # Green circle (outline)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, s, s)
    circle.fill.background()
    circle.line.color.rgb = GREEN
    circle.line.width = Pt(3)

    # Stem (thin tall rectangle)
    stem_w = Pt(5)
    stem_h = Inches(size_inches * 0.28)
    stem_left = Inches(cx_inches) - stem_w // 2
    stem_top = Inches(cy_inches) - Inches(size_inches * 0.22)
    stem = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, stem_left, stem_top, stem_w, stem_h)
    stem.fill.solid()
    stem.fill.fore_color.rgb = WHITE
    stem.line.fill.background()
    stem.rotation = 0

    # Left leaf (small oval, rotated)
    leaf_w = Inches(size_inches * 0.18)
    leaf_h = Inches(size_inches * 0.09)
    leaf_left = Inches(cx_inches) - Inches(size_inches * 0.18)
    leaf_top = Inches(cy_inches) - Inches(size_inches * 0.18)
    left_leaf = slide.shapes.add_shape(MSO_SHAPE.OVAL, leaf_left, leaf_top, leaf_w, leaf_h)
    left_leaf.fill.solid()
    left_leaf.fill.fore_color.rgb = WHITE
    left_leaf.line.fill.background()
    left_leaf.rotation = -40

    # Right leaf
    right_leaf = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx_inches) + Inches(size_inches * 0.02),
        leaf_top,
        leaf_w, leaf_h
    )
    right_leaf.fill.solid()
    right_leaf.fill.fore_color.rgb = WHITE
    right_leaf.line.fill.background()
    right_leaf.rotation = 40

    # Seed (small oval at bottom)
    seed_w = Inches(size_inches * 0.16)
    seed_h = Inches(size_inches * 0.11)
    seed = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx_inches) - seed_w // 2,
        Inches(cy_inches) + Inches(size_inches * 0.10),
        seed_w, seed_h
    )
    seed.fill.solid()
    seed.fill.fore_color.rgb = GREEN
    seed.line.fill.background()


def build_slide_01_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, DARK_BG)
    # Accent line top
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Pt(5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()

    # Logo icon
    draw_logo_icon(slide, cx_inches=6.666, cy_inches=1.8, size_inches=1.4)

    # GRAN + IX wordmark
    add_rich_textbox(
        slide, Inches(1), Inches(2.7), Inches(11.3), Inches(1.2),
        [
            {"text": "GRAN", "size": 72, "color": WHITE, "bold": True, "font": TITLE_FONT},
            {"text": "IX", "size": 72, "color": GREEN, "bold": True, "font": TITLE_FONT},
        ],
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, Inches(1.5), Inches(3.9), Inches(10.3), Inches(0.8),
        "Onde criancas aprendem a fazer grana crescer",
        font_size=26, font_color=LIGHT_GRAY, font_name=BODY_FONT,
        alignment=PP_ALIGN.CENTER,
    )
    # Divider
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.9), Inches(2.3), Pt(2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()

    add_textbox(
        slide, Inches(1.5), Inches(5.3), Inches(10.3), Inches(0.5),
        "Pitch para Parceria Estrategica | 2026",
        font_size=16, font_color=LIGHT_GRAY, font_name=BODY_FONT,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, Inches(1.5), Inches(5.9), Inches(10.3), Inches(0.4),
        "O futuro comeca com um grao",
        font_size=12, font_color=RGBColor(0x66, 0x66, 0x66), font_name=BODY_FONT,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, Inches(9.5), Inches(6.7), Inches(3.5), Inches(0.4),
        "Confidencial", font_size=12, font_color=LIGHT_GRAY,
        font_name=BODY_FONT, alignment=PP_ALIGN.RIGHT,
    )


def build_slide_02_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "O Problema")
    add_green_line(slide, Inches(0.95))

    # Two-column layout: 2 big stats left, 2 big stats right
    stats = [
        ("85%", "dos pais querem ensinar\neducacao financeira aos filhos", GREEN),
        ("42%", "dos adolescentes ainda usam\ncofrinho fisico para poupar", ORANGE),
        ("37,9M", "criancas e adolescentes\nno Brasil (5-17 anos)", BLUE),
        ("93%", "delas ja usam internet\nmas nenhum app resolve o problema", PINK),
    ]

    # 2x2 grid with generous spacing
    box_w = Inches(5.5)
    box_h = Inches(1.7)
    gap_x = Inches(0.5)
    gap_y = Inches(0.4)
    x_positions = [Inches(0.8), Inches(0.8) + box_w + gap_x]
    y_positions = [Inches(1.5), Inches(1.5) + box_h + gap_y]

    for i, (num, lbl, color) in enumerate(stats):
        col = i % 2
        row = i // 2
        left = x_positions[col]
        top = y_positions[row]

        # Card background
        add_rounded_rect(slide, left, top, box_w, box_h, CARD_BG)

        # Colored accent bar on left edge
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, Pt(6), box_h
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # Big number
        add_textbox(
            slide, left + Inches(0.3), top + Inches(0.2),
            Inches(2.0), Inches(0.8),
            num, font_size=44, font_color=color,
            bold=True, font_name=TITLE_FONT, alignment=PP_ALIGN.LEFT,
        )
        # Description
        add_textbox(
            slide, left + Inches(2.4), top + Inches(0.25),
            Inches(2.9), Inches(1.2),
            lbl, font_size=14, font_color=WHITE,
            font_name=BODY_FONT, alignment=PP_ALIGN.LEFT,
        )

    # Quote at bottom with proper spacing
    add_rounded_rect(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.0), DARK_BG2)
    add_textbox(
        slide, Inches(1.2), Inches(5.4), Inches(11.0), Inches(0.8),
        '"Existe um gap enorme entre a consciencia dos pais e as ferramentas disponiveis.\nNos EUA, uma empresa resolveu isso e vale US$ 3,2 bilhoes. No Brasil, ninguem resolveu ainda."',
        font_size=14, font_color=LIGHT_GRAY, font_name=BODY_FONT,
        alignment=PP_ALIGN.CENTER, bold=False,
    )


def build_slide_03_greenlight(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "O Modelo que Funciona: Greenlight (EUA)")
    slide_subtitle(slide, "Fundada em 2014 | Atlanta, GA")
    add_green_line(slide, Inches(0.95))

    metrics = [
        ("6,5M", "familias ativas"),
        ("US$ 228,5M", "receita anual"),
        ("US$ 3,2B", "valuation"),
        ("US$ 556M", "funding total"),
        ("150+", "bancos parceiros"),
        ("US$ 2B", "transacoes/ano\npor criancas"),
    ]
    cols = 3
    box_w = Inches(3.6)
    box_h = Inches(1.6)
    gap_x = Inches(0.35)
    gap_y = Inches(0.3)
    x_start = Inches(0.8)
    y_start = Inches(1.7)
    for i, (num, lbl) in enumerate(metrics):
        r = i // cols
        c = i % cols
        stat_box(
            slide,
            x_start + c * (box_w + gap_x),
            y_start + r * (box_h + gap_y),
            box_w, box_h, num, lbl
        )

    add_textbox(
        slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.5),
        "Investidores: Andreessen Horowitz, JPMorgan Chase, Wells Fargo",
        font_size=14, font_color=LIGHT_GRAY, font_name=BODY_FONT,
        alignment=PP_ALIGN.CENTER,
    )


def build_slide_04_why_greenlight(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "Por que Greenlight Funciona")
    add_green_line(slide, Inches(0.95))

    pillars = [
        (
            "Assinatura Familiar",
            "US$ 5,99-15,98/mes por familia.\n\nIncentivo alinhado: empresa ganha quando familia aprende, nao quando crianca gasta.",
        ),
        (
            "Greenlight for Banks",
            "150+ bancos parceiros incluindo JPMorgan e Wells Fargo.\n\nBancos investem hoje para ter clientes adultos amanha.",
        ),
        (
            "CAC Intergeracional",
            "Bancos gastam US$ 200-400 por cliente adulto.\n\nVia Greenlight, crianca de 10 anos vira cliente aos 18 a custo zero.",
        ),
    ]
    box_w = Inches(3.6)
    box_h = Inches(3.5)
    gap = Inches(0.4)
    x_start = Inches(0.8)
    for i, (title, body) in enumerate(pillars):
        add_card(
            slide, x_start + i * (box_w + gap), Inches(1.5),
            box_w, box_h, title, body,
            title_size=18, body_size=14,
        )


def build_slide_05_market(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "O Mercado Brasileiro")
    add_green_line(slide, Inches(0.95))

    # TAM / SAM / SOM circles (as rounded rects)
    circles = [
        ("TAM", "25,9M familias\ncom filhos < 18", Inches(3.0), BLUE_DEEP),
        ("SAM", "6,7M familias\nclasses A/B/C", Inches(2.5), RGBColor(0x28, 0x30, 0x90)),
        ("SOM (3 anos)", "500.000 familias\n(7,5% do SAM)", Inches(2.0), RGBColor(0x00, 0x60, 0x28)),
    ]
    x = Inches(0.8)
    for label, desc, size, color in circles:
        add_rounded_rect(slide, x, Inches(1.5), size, Inches(2.2), color, border_color=GREEN, border_width=Pt(2))
        add_textbox(
            slide, x + Inches(0.15), Inches(1.65), size - Inches(0.3), Inches(0.5),
            label, font_size=16, font_color=GREEN, bold=True,
            font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, x + Inches(0.15), Inches(2.2), size - Inches(0.3), Inches(1.2),
            desc, font_size=13, font_color=WHITE,
            font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
        )
        x += size + Inches(0.3)

    # Timing points
    timing = [
        "PIX como infraestrutura nativa",
        "90%+ transacoes ja sao mobile",
        "Educacao financeira obrigatoria nas escolas",
        "80% dos teens ja tem conta digital",
        "Nenhum player consolidou a lideranca",
    ]
    y = Inches(4.1)
    for t in timing:
        add_rich_textbox(
            slide, Inches(0.8), y, Inches(11.7), Inches(0.35),
            [
                {"text": "  >  ", "size": 14, "color": GREEN, "bold": True, "font": TITLE_FONT},
                {"text": t, "size": 14, "color": WHITE, "font": BODY_FONT},
            ],
        )
        y += Inches(0.35)


def build_slide_06_competition(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "Panorama Competitivo")
    add_green_line(slide, Inches(0.95))

    headers = ["App", "Cartao", "Tarefas", "Educacao", "Investir", "4 Baldes"]
    rows = [
        ["Nubank Familia", "Sim", "Nao", "Nao", "Nao", "NAO"],
        ["C6 Yellow", "Sim", "Nao", "Nao", "Sim", "NAO"],
        ["NextJoy (Bradesco)", "Sim", "Sim", "Basico", "Nao", "NAO"],
        ["Inter Kids", "Sim", "Nao", "Nao", "Sim", "NAO"],
        ["Mozper", "Sim", "Sim", "Nao", "Nao", "NAO"],
        ["GRANIX", "Sim", "Sim", "Sim", "Sim", "SIM"],
    ]

    col_widths = [Inches(2.5), Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.6)]
    x_start = Inches(0.8)
    y_start = Inches(1.5)
    row_h = Inches(0.52)

    # Header row
    x = x_start
    for j, h in enumerate(headers):
        add_rounded_rect(slide, x, y_start, col_widths[j] - Pt(4), row_h - Pt(4), BLUE_DEEP)
        add_textbox(
            slide, x + Pt(4), y_start + Pt(2), col_widths[j] - Pt(12), row_h - Pt(8),
            h, font_size=12, font_color=GREEN, bold=True,
            font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        x += col_widths[j]

    # Data rows
    for i, row in enumerate(rows):
        is_granix = (i == len(rows) - 1)
        y = y_start + (i + 1) * row_h
        x = x_start
        for j, cell in enumerate(row):
            bg = CARD_HIGHLIGHT if is_granix else (CARD_BG if i % 2 == 0 else DARK_BG2)
            add_rounded_rect(slide, x, y, col_widths[j] - Pt(4), row_h - Pt(4), bg)
            cell_color = GREEN if (cell in ("Sim", "SIM") and (is_granix or j < 5)) else WHITE
            if cell in ("Nao", "NAO", "Basico"):
                cell_color = LIGHT_GRAY
            if is_granix and cell == "SIM":
                cell_color = GREEN
            add_textbox(
                slide, x + Pt(4), y + Pt(2), col_widths[j] - Pt(12), row_h - Pt(8),
                cell, font_size=12, font_color=cell_color,
                bold=is_granix, font_name=BODY_FONT,
                alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            )
            x += col_widths[j]

    add_textbox(
        slide, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.5),
        "Nenhum concorrente combina: cartao + tarefas + educacao + investimentos + 4 Baldes",
        font_size=14, font_color=GREEN, bold=True,
        font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
    )


def build_slide_07_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "GRANIX: O Greenlight Brasileiro")
    add_green_line(slide, Inches(0.95))

    add_textbox(
        slide, Inches(1), Inches(1.3), Inches(11.3), Inches(0.6),
        "GRANA + IX = GRANIX",
        font_size=28, font_color=GREEN, bold=True,
        font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, Inches(1), Inches(1.9), Inches(11.3), Inches(0.5),
        "Plataforma de educacao financeira para familias brasileiras",
        font_size=16, font_color=LIGHT_GRAY, font_name=BODY_FONT,
        alignment=PP_ALIGN.CENTER,
    )

    features = [
        ("Cartao de Debito", "Controle total dos pais, notificacoes em tempo real", BLUE),
        ("4 Baldes Inteligentes", "Gastar, Guardar, Doar, Investir", GREEN),
        ("Tarefas & Recompensas", "Criancas ganham mesada completando tarefas", ORANGE),
        ("GRANIX Academy", "Educacao gamificada por faixa etaria", PINK),
    ]
    box_w = Inches(2.7)
    box_h = Inches(2.5)
    gap = Inches(0.35)
    x_start = Inches(0.8)
    for i, (title, body, accent) in enumerate(features):
        left = x_start + i * (box_w + gap)
        top = Inches(2.8)
        add_rounded_rect(slide, left, top, box_w, box_h, CARD_BG)
        # Accent bar at top of card
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, box_w, Pt(4)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        add_textbox(
            slide, left + Inches(0.2), top + Inches(0.3),
            box_w - Inches(0.4), Inches(0.5),
            title, font_size=16, font_color=accent, bold=True,
            font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, left + Inches(0.2), top + Inches(1.0),
            box_w - Inches(0.4), Inches(1.3),
            body, font_size=13, font_color=WHITE,
            font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
        )


def build_slide_08_buckets(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "O Coracao do Produto: 4 Baldes")
    add_green_line(slide, Inches(0.95))

    buckets = [
        (ORANGE, "GASTAR", "40%", "Saldo disponivel\nno cartao"),
        (BLUE, "GUARDAR", "30%", "Metas de poupanca\nvisuais"),
        (PINK, "DOAR", "10%", "Reserva para\ndoacoes"),
        (GREEN, "INVESTIR", "20%", "CDB, Tesouro Direto\ncom aprovacao dos pais"),
    ]
    box_w = Inches(2.7)
    box_h = Inches(3.0)
    gap = Inches(0.35)
    x_start = Inches(0.8)
    for i, (color, name, pct, desc) in enumerate(buckets):
        left = x_start + i * (box_w + gap)
        top = Inches(1.5)
        add_rounded_rect(slide, left, top, box_w, box_h, CARD_BG, border_color=color, border_width=Pt(3))
        # Color accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, box_w, Pt(6)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        add_textbox(
            slide, left, top + Inches(0.25), box_w, Inches(0.5),
            name, font_size=22, font_color=color, bold=True,
            font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, left, top + Inches(0.8), box_w, Inches(0.6),
            pct, font_size=40, font_color=WHITE, bold=True,
            font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, left + Inches(0.2), top + Inches(1.6), box_w - Inches(0.4), Inches(1.2),
            desc, font_size=13, font_color=LIGHT_GRAY,
            font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
        )

    add_textbox(
        slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(0.5),
        "Cada mesada e automaticamente distribuida. Pais configuram os percentuais.",
        font_size=14, font_color=WHITE, font_name=BODY_FONT,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(0.4),
        "Nenhum concorrente brasileiro tem este sistema.",
        font_size=14, font_color=GREEN, bold=True,
        font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
    )


def build_slide_09_pricing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "Modelo de Monetizacao")
    add_green_line(slide, Inches(0.95))

    plans = [
        ("GRATIS", "R$ 0/mes", "1 filho, cartao virtual,\n5 tarefas/mes", CARD_BG, None),
        ("FAMILIA", "R$ 29/mes", "4 filhos, cartao fisico,\ntarefas ilimitadas,\nAcademy basico", CARD_BG, GREEN),
        ("FAMILIA+", "R$ 49/mes", "6 filhos, todas as trilhas,\ninvestimentos, cashback 1%,\nrelatorios IA", CARD_BG, None),
    ]
    box_w = Inches(3.6)
    box_h = Inches(3.5)
    gap = Inches(0.35)
    x_start = Inches(0.8)
    for i, (name, price, features, bg, highlight) in enumerate(plans):
        left = x_start + i * (box_w + gap)
        top = Inches(1.5)
        border = highlight if highlight else None
        add_rounded_rect(slide, left, top, box_w, box_h, bg, border_color=border, border_width=Pt(3) if border else Pt(0))

        if highlight:
            tag = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left + box_w - Inches(1.5), top - Inches(0.2), Inches(1.5), Inches(0.35),
            )
            tag.fill.solid()
            tag.fill.fore_color.rgb = GREEN
            tag.line.fill.background()
            add_textbox(
                slide, left + box_w - Inches(1.5), top - Inches(0.18), Inches(1.5), Inches(0.3),
                "POPULAR", font_size=10, font_color=WHITE, bold=True,
                font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
            )

        add_textbox(
            slide, left, top + Inches(0.3), box_w, Inches(0.4),
            name, font_size=20, font_color=GREEN, bold=True,
            font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, left, top + Inches(0.8), box_w, Inches(0.5),
            price, font_size=28, font_color=WHITE, bold=True,
            font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, left + Inches(0.3), top + Inches(1.5), box_w - Inches(0.6), Inches(1.8),
            features, font_size=13, font_color=LIGHT_GRAY,
            font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
        )

    # Revenue split
    add_textbox(
        slide, Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.5),
        "Assinaturas 60%  |  Interchange 25%  |  Float 10%  |  Investimentos 5%",
        font_size=15, font_color=GREEN, bold=True,
        font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
    )


def build_slide_10_academy(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "GRANIX Academy")
    slide_subtitle(slide, "Educacao financeira gamificada por idade")
    add_green_line(slide, Inches(0.95))

    tracks = [
        ("Semente (6-8 anos)", "8 licoes\nVideos curtos + HQs com mascote", RGBColor(0x66, 0xBB, 0x6A)),
        ("Broto (9-11 anos)", "12 licoes\nQuizzes + simuladores", BLUE),
        ("Arvore (12-14 anos)", "16 licoes\nCases reais + desafios praticos", ORANGE),
        ("Floresta (15-17 anos)", "20 licoes\nInvestimentos + empreendedorismo", GREEN),
    ]
    box_w = Inches(2.7)
    box_h = Inches(2.5)
    gap = Inches(0.35)
    x_start = Inches(0.8)
    for i, (title, body, accent) in enumerate(tracks):
        left = x_start + i * (box_w + gap)
        top = Inches(1.7)
        add_rounded_rect(slide, left, top, box_w, box_h, CARD_BG)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, box_w, Pt(4)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        add_textbox(
            slide, left + Inches(0.15), top + Inches(0.3),
            box_w - Inches(0.3), Inches(0.5),
            title, font_size=15, font_color=accent, bold=True,
            font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, left + Inches(0.15), top + Inches(1.0),
            box_w - Inches(0.3), Inches(1.3),
            body, font_size=13, font_color=WHITE,
            font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
        )

    add_textbox(
        slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(0.5),
        "Recompensa real: R$ 5 por trilha + R$ 20 bonus por completar todas",
        font_size=15, font_color=GREEN, bold=True,
        font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
    )


def build_slide_11_why_xp(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "Por que XP / Banco Modal")
    slide_subtitle(slide, "Voces tem o que precisamos. Nos temos o que voces precisam.")
    add_green_line(slide, Inches(0.95))

    boxes = [
        ("4,7M", "clientes ativos", "Adultos classes A/B, a maioria pais. Base pronta."),
        ("18.200", "assessores", "Canal de venda consultiva de alta confianca."),
        ("R$ 1,2T", "em AUM", "Clientes com poder de compra para plano premium."),
        ("Instituto XP", "", "Meta: educacao financeira para 50M de brasileiros."),
    ]
    box_w = Inches(5.5)
    box_h = Inches(1.1)
    gap = Inches(0.2)
    x_start = Inches(0.8)
    y_start = Inches(1.7)
    for i, (num, subtitle, desc) in enumerate(boxes):
        r = i // 2
        c = i % 2
        left = x_start + c * (box_w + gap)
        top = y_start + r * (box_h + gap)
        add_rounded_rect(slide, left, top, box_w, box_h, CARD_BG)
        add_textbox(
            slide, left + Inches(0.2), top + Inches(0.1),
            Inches(1.8), Inches(0.5),
            num, font_size=24, font_color=GREEN, bold=True,
            font_name=TITLE_FONT,
        )
        if subtitle:
            add_textbox(
                slide, left + Inches(0.2), top + Inches(0.55),
                Inches(1.8), Inches(0.35),
                subtitle, font_size=11, font_color=LIGHT_GRAY,
                font_name=BODY_FONT,
            )
        add_textbox(
            slide, left + Inches(2.1), top + Inches(0.15),
            Inches(3.2), Inches(0.8),
            desc, font_size=13, font_color=WHITE,
            font_name=BODY_FONT,
        )

    # Quote
    add_textbox(
        slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.5),
        '"Se 5% dos clientes XP aderirem, sao 235.000 familias."',
        font_size=18, font_color=GREEN, bold=True,
        font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
    )


def build_slide_12_cac(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "A Logica Estrategica: CAC Intergeracional")
    add_green_line(slide, Inches(0.95))

    comparisons = [
        ("R$ 500-1.000", "Custo para adquirir\n1 cliente adulto hoje", RED),
        ("~R$ 50", "Custo para onboardar\n1 familia via GRANIX", YELLOW),
        ("R$ 0", "Custo quando crianca\nvira cliente XP aos 18", GREEN),
    ]
    box_w = Inches(3.6)
    box_h = Inches(2.8)
    gap = Inches(0.35)
    x_start = Inches(0.8)
    for i, (amount, desc, color) in enumerate(comparisons):
        left = x_start + i * (box_w + gap)
        top = Inches(1.5)
        add_rounded_rect(slide, left, top, box_w, box_h, CARD_BG, border_color=color, border_width=Pt(3))
        add_textbox(
            slide, left, top + Inches(0.4), box_w, Inches(0.8),
            amount, font_size=36, font_color=color, bold=True,
            font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, left + Inches(0.3), top + Inches(1.4), box_w - Inches(0.6), Inches(1.2),
            desc, font_size=14, font_color=WHITE,
            font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
        )

    add_textbox(
        slide, Inches(0.8), Inches(4.7), Inches(11.7), Inches(0.5),
        "Cada familia GRANIX = pipeline de 1-4 futuros clientes adultos da XP",
        font_size=16, font_color=WHITE, bold=True,
        font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.4),
        "Greenlight provou isso com 150+ bancos nos EUA, incluindo JPMorgan",
        font_size=14, font_color=LIGHT_GRAY,
        font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
    )


def build_slide_13_partnership(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "Opcoes de Parceria")
    add_green_line(slide, Inches(0.95))

    options = [
        (
            "Opcao A: BaaS + Distribuicao",
            "XP fornece infra bancaria (Modal) + canal assessores.\n\nRevenue share sobre assinaturas.",
            None,
        ),
        (
            "Opcao B: BaaS + Investimento",
            "Opcao A + XP investe R$ 1-2M no pre-seed.\n\nEquity minoritario.",
            GREEN,
        ),
        (
            "Opcao C: White-Label",
            "GRANIX como 'XP Familia'.\nXP paga licenciamento.\n\nModelo Greenlight for Banks.",
            None,
        ),
    ]
    box_w = Inches(3.6)
    box_h = Inches(3.5)
    gap = Inches(0.35)
    x_start = Inches(0.8)
    for i, (title, body, highlight) in enumerate(options):
        left = x_start + i * (box_w + gap)
        top = Inches(1.5)
        border = highlight if highlight else None
        add_rounded_rect(slide, left, top, box_w, box_h, CARD_BG, border_color=border, border_width=Pt(3) if border else Pt(0))

        if highlight:
            tag = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left + box_w - Inches(2.0), top - Inches(0.22), Inches(2.0), Inches(0.35),
            )
            tag.fill.solid()
            tag.fill.fore_color.rgb = GREEN
            tag.line.fill.background()
            add_textbox(
                slide, left + box_w - Inches(2.0), top - Inches(0.20), Inches(2.0), Inches(0.3),
                "RECOMENDADO", font_size=10, font_color=WHITE, bold=True,
                font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
            )

        add_textbox(
            slide, left + Inches(0.25), top + Inches(0.3),
            box_w - Inches(0.5), Inches(0.5),
            title, font_size=16, font_color=GREEN, bold=True,
            font_name=TITLE_FONT, alignment=PP_ALIGN.LEFT,
        )
        add_textbox(
            slide, left + Inches(0.25), top + Inches(1.0),
            box_w - Inches(0.5), Inches(2.3),
            body, font_size=13, font_color=WHITE,
            font_name=BODY_FONT, alignment=PP_ALIGN.LEFT,
        )


def build_slide_14_investment(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "Investimento e Timeline")
    add_green_line(slide, Inches(0.95))

    # Left side - Budget
    add_rounded_rect(slide, Inches(0.8), Inches(1.5), Inches(5.2), Inches(1.6), CARD_BG)
    add_textbox(
        slide, Inches(1.0), Inches(1.6), Inches(4.8), Inches(0.4),
        "Investimento", font_size=18, font_color=GREEN, bold=True,
        font_name=TITLE_FONT,
    )
    add_textbox(
        slide, Inches(1.0), Inches(2.1), Inches(4.8), Inches(0.35),
        "MVP + 6 meses: R$ 1,1M", font_size=14, font_color=WHITE,
        font_name=BODY_FONT,
    )
    add_textbox(
        slide, Inches(1.0), Inches(2.45), Inches(4.8), Inches(0.35),
        "Com parceria XP: ~R$ 800-900K", font_size=14, font_color=WHITE,
        font_name=BODY_FONT,
    )

    # Right side - Timeline
    add_rounded_rect(slide, Inches(6.3), Inches(1.5), Inches(6.2), Inches(3.2), CARD_BG)
    add_textbox(
        slide, Inches(6.5), Inches(1.6), Inches(5.8), Inches(0.4),
        "Timeline", font_size=18, font_color=GREEN, bold=True,
        font_name=TITLE_FONT,
    )
    timeline = [
        "M1-M2: Setup e integracao BaaS",
        "M3-M5: Desenvolvimento MVP",
        "M6-M7: Piloto com 500 familias (assessores XP)",
        "M8: Analise de metricas",
        "M9+: Lancamento publico",
    ]
    y = Inches(2.1)
    for t in timeline:
        add_rich_textbox(
            slide, Inches(6.5), y, Inches(5.8), Inches(0.35),
            [
                {"text": "  >  ", "size": 12, "color": GREEN, "bold": True, "font": TITLE_FONT},
                {"text": t, "size": 12, "color": WHITE, "font": BODY_FONT},
            ],
        )
        y += Inches(0.38)

    # Projections table
    add_textbox(
        slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.4),
        "Projecoes Financeiras", font_size=18, font_color=GREEN, bold=True,
        font_name=TITLE_FONT, alignment=PP_ALIGN.LEFT,
    )

    proj_headers = ["", "Familias", "ARR", "EBITDA"]
    proj_rows = [
        ["Ano 1", "15K", "R$ 4,1M", "-"],
        ["Ano 2", "45K", "R$ 12,2M", "-"],
        ["Ano 3", "100K", "R$ 27M", "-"],
        ["Ano 5", "300K", "R$ 81,1M", "51%"],
    ]
    col_w = [Inches(1.5), Inches(2.0), Inches(2.0), Inches(2.0)]
    x_start = Inches(1.8)
    row_h = Inches(0.38)
    y_start = Inches(5.5)

    # Header
    x = x_start
    for j, h in enumerate(proj_headers):
        add_textbox(
            slide, x, y_start, col_w[j], row_h,
            h, font_size=11, font_color=GREEN, bold=True,
            font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
        )
        x += col_w[j]

    for i, row in enumerate(proj_rows):
        x = x_start
        y = y_start + (i + 1) * row_h
        for j, cell in enumerate(row):
            color = GREEN if j == 0 else WHITE
            add_textbox(
                slide, x, y, col_w[j], row_h,
                cell, font_size=11, font_color=color,
                bold=(j == 0), font_name=BODY_FONT,
                alignment=PP_ALIGN.CENTER,
            )
            x += col_w[j]


def build_slide_15_cta(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "O Ask")
    add_green_line(slide, Inches(0.95))

    asks = [
        ("1. Infraestrutura bancaria", "Usar a plataforma do Modal como BaaS"),
        ("2. Canal de distribuicao", "Acesso aos assessores e clientes XP para o piloto"),
        ("3. Investimento pre-seed", "R$ 1,5M por equity minoritario"),
    ]
    box_w = Inches(3.6)
    box_h = Inches(1.8)
    gap = Inches(0.35)
    x_start = Inches(0.8)
    for i, (title, body) in enumerate(asks):
        left = x_start + i * (box_w + gap)
        top = Inches(1.5)
        add_rounded_rect(slide, left, top, box_w, box_h, CARD_BG, border_color=GREEN, border_width=Pt(2))
        add_textbox(
            slide, left + Inches(0.2), top + Inches(0.25),
            box_w - Inches(0.4), Inches(0.5),
            title, font_size=16, font_color=GREEN, bold=True,
            font_name=TITLE_FONT, alignment=PP_ALIGN.LEFT,
        )
        add_textbox(
            slide, left + Inches(0.2), top + Inches(0.85),
            box_w - Inches(0.4), Inches(0.8),
            body, font_size=13, font_color=WHITE,
            font_name=BODY_FONT, alignment=PP_ALIGN.LEFT,
        )

    # Big closing quote
    add_rounded_rect(slide, Inches(0.8), Inches(3.8), Inches(11.7), Inches(1.5), CARD_BG)
    add_textbox(
        slide, Inches(1.2), Inches(3.9), Inches(11.0), Inches(1.3),
        '"A questao nao e SE alguem vai fazer o Greenlight brasileiro.\nE QUEM e QUANDO."',
        font_size=22, font_color=WHITE, bold=True,
        font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.5),
        "parceiros@granix.app",
        font_size=18, font_color=GREEN, bold=False,
        font_name=BODY_FONT, alignment=PP_ALIGN.CENTER,
    )

    # Confidential
    add_textbox(
        slide, Inches(9.5), Inches(6.5), Inches(3.5), Inches(0.4),
        "Confidencial", font_size=12, font_color=LIGHT_GRAY,
        font_name=BODY_FONT, alignment=PP_ALIGN.RIGHT,
    )


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    template_path = (
        "/Users/lab/Claude/templates-powerpoint/allpowerpointtemplates/"
        "pitch-deck-powerpoint/01 Presentation/"
        "Pitch Deck PowerPoint Template Dark.pptx"
    )
    output_path = "/Users/lab/Claude/Projetos/Granix-App/docs/GRANIX-Pitch-XP-Modal.pptx"

    # Load template for its theme/master, then remove existing slides
    prs = Presentation(template_path)

    # Set 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Remove all existing template slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(qn("r:id"))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # Build all 15 slides
    build_slide_01_cover(prs)
    build_slide_02_problem(prs)
    build_slide_03_greenlight(prs)
    build_slide_04_why_greenlight(prs)
    build_slide_05_market(prs)
    build_slide_06_competition(prs)
    build_slide_07_solution(prs)
    build_slide_08_buckets(prs)
    build_slide_09_pricing(prs)
    build_slide_10_academy(prs)
    build_slide_11_why_xp(prs)
    build_slide_12_cac(prs)
    build_slide_13_partnership(prs)
    build_slide_14_investment(prs)
    build_slide_15_cta(prs)

    prs.save(output_path)
    print(f"Saved {len(prs.slides)} slides to: {output_path}")


if __name__ == "__main__":
    main()
