"""
GRANIX Partner Pitch Deck Generator v2 — Component Library + AI Engine.

Generates a 15-slide professional PowerPoint for institutional partners.
Focus: product flow, content, AI differentiation, white-label integration.
Dark theme with GRANIX brand colors and 16:9 aspect ratio.

Usage:
    python generate_pitch_partners.py
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Constants & Brand Identity
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

GREEN = RGBColor(0x00, 0xC8, 0x53)
BLUE_DEEP = RGBColor(0x1A, 0x23, 0x7E)
DARK_BG = RGBColor(0x0F, 0x12, 0x20)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xFF, 0x6D, 0x00)
PINK = RGBColor(0xE9, 0x1E, 0x63)
LIGHT_GRAY = RGBColor(0xB0, 0xB0, 0xB0)
MED_GRAY = RGBColor(0x70, 0x70, 0x80)
CARD_BG = RGBColor(0x1A, 0x1E, 0x30)
CARD_BG_ALT = RGBColor(0x22, 0x27, 0x3A)
PURPLE = RGBColor(0x9C, 0x27, 0xB0)
RED = RGBColor(0xF4, 0x43, 0x36)
YELLOW = RGBColor(0xFF, 0xEB, 0x3B)
BLUE = RGBColor(0x21, 0x96, 0xF3)
TEAL = RGBColor(0x00, 0xBF, 0xA5)

TITLE_FONT = "Poppins"
BODY_FONT = "Inter"

OUTPUT_PATH = "/Users/lab/Claude/Projetos/Granix-App/docs/GRANIX-Pitch-Partners-v2.pptx"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def set_slide_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(
    slide, left, top, width, height, text,
    font_size=18, font_color=WHITE, bold=False,
    font_name=BODY_FONT, alignment=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = alignment
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    txBox.text_frame._txBody.bodyPr.set(
        "anchor",
        {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}.get(anchor, "t"),
    )
    return txBox


def add_multiline_textbox(
    slide, left, top, width, height,
    lines: list[tuple[str, int, RGBColor, bool, str]],
    alignment=PP_ALIGN.LEFT, line_spacing=1.2,
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, (text, size, color, bld, fname) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bld
        p.font.name = fname
        p.alignment = alignment
        p.space_after = Pt(size * (line_spacing - 1.0))
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title_text, body_text,
             fill=CARD_BG, title_color=GREEN, body_color=WHITE,
             title_size=16, body_size=13, border_color=None, accent_color=None):
    add_rounded_rect(slide, left, top, width, height, fill,
                     border_color=border_color, border_width=Pt(2) if border_color else Pt(0))
    if accent_color:
        add_rect(slide, left, top, width, Pt(5), accent_color)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.2),
                width - Inches(0.4), Inches(0.4),
                title_text, font_size=title_size, font_color=title_color,
                bold=True, font_name=TITLE_FONT)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.65),
                width - Inches(0.4), height - Inches(0.85),
                body_text, font_size=body_size, font_color=body_color,
                font_name=BODY_FONT)


def slide_title(slide, text, top=Inches(0.3), font_size=32):
    add_textbox(slide, Inches(0.8), top, Inches(11.7), Inches(0.8),
                text, font_size=font_size, font_color=WHITE, bold=True, font_name=TITLE_FONT)


def slide_subtitle(slide, text, top=Inches(1.1)):
    add_textbox(slide, Inches(0.8), top, Inches(11.7), Inches(0.5),
                text, font_size=18, font_color=LIGHT_GRAY, font_name=BODY_FONT)


def slide_number_badge(slide, number):
    add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.35),
                number, font_size=10, font_color=MED_GRAY,
                alignment=PP_ALIGN.RIGHT, font_name=BODY_FONT)


def add_stat_card(slide, left, top, number, label, color=GREEN):
    w, h = Inches(2.6), Inches(1.5)
    add_rounded_rect(slide, left, top, w, h, CARD_BG, border_color=color, border_width=Pt(1))
    add_textbox(slide, left, top + Inches(0.15), w, Inches(0.7),
                number, font_size=36, font_color=color, bold=True,
                font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + Inches(0.9), w, Inches(0.4),
                label, font_size=12, font_color=LIGHT_GRAY,
                font_name=BODY_FONT, alignment=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide Builders
# ---------------------------------------------------------------------------
def build_slide_01_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    # Accent bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Pt(4), GREEN)
    # Logo / Title
    add_multiline_textbox(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(3.5), [
        ("GRANIX", 72, GREEN, True, TITLE_FONT),
        ("", 12, DARK_BG, False, BODY_FONT),
        ("Component Library + AI Content Engine", 28, WHITE, False, TITLE_FONT),
        ("para Educacao Financeira de Jovens", 28, WHITE, False, TITLE_FONT),
        ("", 12, DARK_BG, False, BODY_FONT),
        ("Embeddavel  •  White-label  •  AI-powered  •  Zero integracao com core bancario", 16, LIGHT_GRAY, False, BODY_FONT),
    ], alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(3), Inches(6.2), Inches(7), Inches(0.5),
                "Pitch para Parceiros Institucionais  |  Marco 2026",
                font_size=14, font_color=MED_GRAY, alignment=PP_ALIGN.CENTER)


def build_slide_02_gap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "O Gap: Conta de Menor Sem Experiencia")
    slide_subtitle(slide, "Sua instituicao ja tem a infraestrutura. Falta a experiencia que engaja.")
    slide_number_badge(slide, "02")

    # What you HAVE
    add_card(slide, Inches(0.6), Inches(1.8), Inches(5.5), Inches(5.0),
             "O que voce JA TEM", (
                 "✅  Conta transacional (PIX, cartao)\n"
                 "✅  Investimentos para menores\n"
                 "✅  Compliance e regulacao\n"
                 "✅  Base de clientes com filhos\n"
                 "✅  App funcional\n"
                 "✅  Assessores como canal"
             ), accent_color=GREEN, title_size=20, body_size=16)

    # What you DON'T HAVE
    add_card(slide, Inches(7.2), Inches(1.8), Inches(5.5), Inches(5.0),
             "O que voce NAO TEM", (
                 "❌  Experiencia adaptada para jovens\n"
                 "❌  Educacao financeira interativa\n"
                 "❌  Gamificacao de qualquer tipo\n"
                 "❌  Sistema de tarefas/recompensas\n"
                 "❌  AI personalizada por perfil\n"
                 "❌  Framework financeiro (4 Potes)"
             ), accent_color=RED, title_color=RED, title_size=20, body_size=16)


def build_slide_03_barrier(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "A Barreira: Nenhuma API Externa no Ambiente Seguro")
    slide_number_badge(slide, "03")

    add_multiline_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.0), [
        ("Nenhuma instituicao financeira vai abrir seu ambiente seguro", 22, WHITE, False, BODY_FONT),
        ("de transacao para uma API externa.", 22, WHITE, False, BODY_FONT),
        ("", 10, DARK_BG, False, BODY_FONT),
        ("Por isso, GRANIX nao integra — GRANIX embeda.", 24, GREEN, True, TITLE_FONT),
    ])

    # Architecture diagram as cards
    add_card(slide, Inches(0.6), Inches(3.8), Inches(5.8), Inches(3.2),
             "DENTRO DO SEU APP (ambiente seguro)", (
                 "GRANIX Component Library\n"
                 "• Componentes React/RN que voce importa\n"
                 "• Recebem dados via props\n"
                 "• Emitem callbacks (voce processa)\n"
                 "• Zero acesso a APIs financeiras\n"
                 "• Themeable com sua marca"
             ), accent_color=GREEN, body_size=14)

    add_card(slide, Inches(7.0), Inches(3.8), Inches(5.8), Inches(3.2),
             "GRANIX CLOUD (externo, sem PII)", (
                 "AI Engine + Content CDN\n"
                 "• Personaliza conteudo por perfil\n"
                 "• Prediz churn e comportamento\n"
                 "• Gera conteudo automaticamente\n"
                 "• Recebe APENAS dados anonimizados\n"
                 "• Zero dados financeiros"
             ), accent_color=BLUE, title_color=BLUE, body_size=14)


def build_slide_04_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "GRANIX: 4 Pilares em 1 Library")
    slide_subtitle(slide, "Tudo que falta na experiencia do seu app jovem, pronto para embeddar.")
    slide_number_badge(slide, "04")

    cards = [
        ("Tarefas & Recompensas", "Pais criam tarefas, filhos completam.\nMesada digital, streaks, aprovacao.\nMotor de engajamento familiar.", GREEN),
        ("GRANIX Academy", "145 licoes por faixa etaria (6-17).\nAlinhado a BNCC.\nQuizzes, simuladores, desafios.", BLUE),
        ("Gamificacao", "XP, badges, niveis, avatar, ranking.\nDesafios semanais personalizados.\nGamificacao alimentada por AI.", PURPLE),
        ("4 Potes", "Gastar / Guardar / Doar / Investir.\nFramework visual educativo.\nZero transacao — voce controla.", ORANGE),
    ]
    for i, (title, body, color) in enumerate(cards):
        left = Inches(0.5 + i * 3.15)
        add_card(slide, left, Inches(2.0), Inches(2.95), Inches(4.8),
                 title, body, accent_color=color, title_color=color,
                 title_size=16, body_size=13)


def build_slide_05_onboarding(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "Fluxo: Onboarding em 3 Minutos")
    slide_subtitle(slide, "5 telas. O sistema se molda ao usuario. Momento 'aha!' no final.")
    slide_number_badge(slide, "05")

    steps = [
        ("1", "Familia", "30s", "Quantos filhos?\nIdades?", GREEN),
        ("2", "Perfil Pai", "30s", "Relacao com dinheiro\nPreocupacoes", BLUE),
        ("3", "Perfil Jovem", "60s", "Interesses\nR$100: o que faria?", ORANGE),
        ("4", "Quiz Flash", "60s", "5 perguntas rapidas\nCalibra nivel", PURPLE),
        ("5", "Momento Aha!", "30s", "Trilha montada\n1o badge + 4 Potes", GREEN),
    ]
    for i, (num, title, time, desc, color) in enumerate(steps):
        left = Inches(0.4 + i * 2.55)
        top = Inches(2.2)
        w = Inches(2.35)
        add_rounded_rect(slide, left, top, w, Inches(4.5), CARD_BG, border_color=color, border_width=Pt(1))
        # Number circle
        add_rounded_rect(slide, left + Inches(0.85), top + Inches(0.2), Inches(0.6), Inches(0.6), color)
        add_textbox(slide, left + Inches(0.85), top + Inches(0.2), Inches(0.6), Inches(0.6),
                    num, font_size=24, font_color=DARK_BG, bold=True,
                    font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Title
        add_textbox(slide, left + Inches(0.1), top + Inches(1.0), w - Inches(0.2), Inches(0.4),
                    title, font_size=16, font_color=color, bold=True,
                    font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER)
        # Time
        add_textbox(slide, left + Inches(0.1), top + Inches(1.45), w - Inches(0.2), Inches(0.3),
                    time, font_size=12, font_color=MED_GRAY,
                    font_name=BODY_FONT, alignment=PP_ALIGN.CENTER)
        # Description
        add_textbox(slide, left + Inches(0.15), top + Inches(1.9), w - Inches(0.3), Inches(2.2),
                    desc, font_size=13, font_color=WHITE,
                    font_name=BODY_FONT, alignment=PP_ALIGN.CENTER)

    # AI note
    add_textbox(slide, Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4),
                "AI processa cada resposta em tempo real — a trilha e personalizada ANTES do usuario comecar.",
                font_size=13, font_color=GREEN, font_name=BODY_FONT, alignment=PP_ALIGN.CENTER)


def build_slide_06_first_session(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "Fluxo: Primeira Sessao — Hook Imediato")
    slide_subtitle(slide, "5-7 minutos. O jovem sai com licao, badge, quiz e objetivo nos 4 Potes.")
    slide_number_badge(slide, "06")

    # What happens
    items = [
        ("1 licao contextualizada", "Tema conectado ao interesse\n(games, futebol, musica)", GREEN),
        ("1 quiz aprovado", "Feedback instantaneo\n'Eu sei coisas!'", BLUE),
        ("1 badge conquistado", "Dopamina imediata\n'Quero mais!'", YELLOW),
        ("1 objetivo nos 4 Potes", "Meta pessoal definida\n'Tenho um plano!'", ORANGE),
    ]
    for i, (title, desc, color) in enumerate(items):
        left = Inches(0.5 + i * 3.15)
        add_card(slide, left, Inches(2.0), Inches(2.95), Inches(2.8),
                 title, desc, accent_color=color, title_color=color,
                 title_size=15, body_size=13)

    # Contextualization examples
    add_textbox(slide, Inches(0.8), Inches(5.1), Inches(11.7), Inches(0.4),
                "Primeira licao e SEMPRE contextualizada ao interesse do jovem:",
                font_size=14, font_color=WHITE, bold=True, font_name=TITLE_FONT)

    examples = [
        ("Games:", "\"Quanto custa um jogo? O preco real das skins\""),
        ("Futebol:", "\"O salario dos jogadores: de onde vem o dinheiro?\""),
        ("Musica:", "\"Como artistas ganham dinheiro com streaming?\""),
    ]
    for i, (ctx, ex) in enumerate(examples):
        add_textbox(slide, Inches(0.8 + i * 4.0), Inches(5.6), Inches(3.8), Inches(0.8),
                    f"{ctx}  {ex}", font_size=12, font_color=LIGHT_GRAY, font_name=BODY_FONT)

    # Parent notification
    add_textbox(slide, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.4),
                "Pai recebe: \"Joao comecou sua jornada financeira! Crie a primeira tarefa para ele.\"",
                font_size=13, font_color=GREEN, font_name=BODY_FONT, alignment=PP_ALIGN.CENTER)


def build_slide_07_weekly_loop(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "Fluxo: Loop Semanal de Engajamento")
    slide_subtitle(slide, "AI gerencia o ritmo. Cada dia tem um proposito. Familia engajada junta.")
    slide_number_badge(slide, "07")

    days = [
        ("SEG", "Desafio da\nsemana", "AI-personalizado", GREEN),
        ("TER-QUI", "1-2 licoes\ncurtas", "5-8 min cada", BLUE),
        ("SEX", "Quiz semanal\n+ XP", "Revisao + badge", PURPLE),
        ("SAB", "Desafio\npratico", "Mundo real!", ORANGE),
        ("DOM", "Family time\n+ resumo", "Pai + filho", PINK),
    ]
    for i, (day, activity, detail, color) in enumerate(days):
        left = Inches(0.3 + i * 2.55)
        top = Inches(2.2)
        w = Inches(2.35)
        add_rounded_rect(slide, left, top, w, Inches(3.5), CARD_BG, border_color=color, border_width=Pt(1))
        add_textbox(slide, left, top + Inches(0.2), w, Inches(0.4),
                    day, font_size=18, font_color=color, bold=True,
                    font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, left, top + Inches(0.8), w, Inches(1.2),
                    activity, font_size=15, font_color=WHITE,
                    font_name=BODY_FONT, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, left, top + Inches(2.3), w, Inches(0.6),
                    detail, font_size=12, font_color=MED_GRAY,
                    font_name=BODY_FONT, alignment=PP_ALIGN.CENTER)

    # Cadence table
    add_textbox(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.4),
                "Cadencia por faixa:  Semente (6-8): 3-4x/sem, 3-5min  |  Broto (9-11): 4-5x, 5-7min  |  Arvore (12-14): 4-5x, 7-10min  |  Floresta (15-17): 3-5x, 8-12min",
                font_size=12, font_color=LIGHT_GRAY, font_name=BODY_FONT, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.4),
                "AI ajusta dinamicamente: mais engajamento → mais conteudo. Menos → reduz para nao saturar.",
                font_size=13, font_color=GREEN, font_name=BODY_FONT, alignment=PP_ALIGN.CENTER)


def build_slide_08_content(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "Conteudo: 145 Licoes, 4 Trilhas, Alinhado a BNCC")
    slide_subtitle(slide, "Curriculo financeiro completo dos 6 aos 17 anos. AI cria e adapta continuamente.")
    slide_number_badge(slide, "08")

    trails = [
        ("Semente", "6-8 anos", "25 licoes", "Dinheiro, troca,\ncofrinho, mesada", GREEN, "Descobrindo\no Dinheiro"),
        ("Broto", "9-11 anos", "33 licoes", "Orcamento, metas,\nempreendedorismo", BLUE, "Construindo\nObjetivos"),
        ("Arvore", "12-14 anos", "43 licoes", "Juros compostos,\ninvestimentos, credito", ORANGE, "Investindo\nno Futuro"),
        ("Floresta", "15-17 anos", "44 licoes", "Bolsa, previdencia,\nplanejamento de vida", PURPLE, "Dominando\no Mercado"),
    ]
    for i, (name, age, lessons, topics, color, subtitle) in enumerate(trails):
        left = Inches(0.5 + i * 3.15)
        top = Inches(2.0)
        w = Inches(2.95)
        h = Inches(4.2)
        add_rounded_rect(slide, left, top, w, h, CARD_BG, border_color=color, border_width=Pt(1))
        # Trail name
        add_textbox(slide, left, top + Inches(0.2), w, Inches(0.5),
                    name, font_size=22, font_color=color, bold=True,
                    font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER)
        # Age
        add_textbox(slide, left, top + Inches(0.7), w, Inches(0.3),
                    age, font_size=14, font_color=MED_GRAY,
                    font_name=BODY_FONT, alignment=PP_ALIGN.CENTER)
        # Lessons count
        add_textbox(slide, left, top + Inches(1.05), w, Inches(0.4),
                    lessons, font_size=16, font_color=WHITE, bold=True,
                    font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER)
        # Subtitle
        add_textbox(slide, left, top + Inches(1.5), w, Inches(0.6),
                    f"\"{subtitle}\"", font_size=12, font_color=LIGHT_GRAY,
                    font_name=BODY_FONT, alignment=PP_ALIGN.CENTER)
        # Topics
        add_textbox(slide, left + Inches(0.15), top + Inches(2.3), w - Inches(0.3), Inches(1.5),
                    topics, font_size=13, font_color=WHITE,
                    font_name=BODY_FONT, alignment=PP_ALIGN.CENTER)

    # Formats
    add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
                "8 formatos: Historias Animadas  •  Quizzes  •  Simuladores  •  Desafios Praticos  •  Videos Curtos  •  Cenarios Interativos  •  Infograficos  •  Flashcards",
                font_size=12, font_color=LIGHT_GRAY, font_name=BODY_FONT, alignment=PP_ALIGN.CENTER)


def build_slide_09_ai_personalization(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "AI: Personalizacao que Transforma 145 Licoes em Milhoes de Jornadas")
    slide_number_badge(slide, "09")

    # Learning Profile card
    add_card(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5),
             "Learning Profile (por usuario)", (
                 "A AI constroi um perfil unico para cada jovem:\n\n"
                 "• Faixa etaria + meses na plataforma\n"
                 "• Formato preferido (simulador, quiz, video)\n"
                 "• Duracao ideal de sessao (ex: 8 min)\n"
                 "• Horario pico de engajamento\n"
                 "• Areas fortes e fracas\n"
                 "• Estilo de erro (impulsivo, cauteloso)\n"
                 "• Interesses (futebol, games, musica)\n"
                 "• Engajamento do pai\n\n"
                 "Cada dado refina a experiencia em tempo real."
             ), accent_color=BLUE, title_color=BLUE, body_size=14)

    # What AI decides
    add_card(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.5),
             "O que a AI decide", (
                 "• Proxima licao (tema + formato + dificuldade)\n"
                 "• Quando enviar notificacao (horario otimo)\n"
                 "• Qual contexto usar (futebol? games?)\n"
                 "• Quando avancar de nivel\n"
                 "• Quando intervir (frustacao ou tedio)"
             ), accent_color=GREEN, body_size=14)

    # Micro-segments
    add_card(slide, Inches(6.8), Inches(4.3), Inches(6.0), Inches(2.7),
             "6 Micro-Segmentos AI-Detected", (
                 "Explorador — alta curiosidade, pula entre temas\n"
                 "Focado — aprofunda um tema ate dominar\n"
                 "Social — engaja com rankings e competicao\n"
                 "Pratico — prefere simuladores, hands-on\n"
                 "Cauteloso — progride devagar, precisa reforco\n"
                 "Intermitente — usa em rajadas, sessoes curtas"
             ), accent_color=PURPLE, title_color=PURPLE, body_size=13)


def build_slide_10_ai_prediction(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "AI: Predicao de Churn + Criacao Automatizada de Conteudo")
    slide_number_badge(slide, "10")

    # Churn prediction
    add_card(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.5),
             "Predicao de Churn em 3 Camadas", (
                 "Alerta Amarelo (7-14 dias antes)\n"
                 "• Reducao de sessoes, quizzes ignorados\n"
                 "→ AI ajusta dificuldade + desafio personalizado\n\n"
                 "Alerta Laranja (3-7 dias antes)\n"
                 "• Zero sessoes por 3+ dias, pai inativo\n"
                 "→ 'Missao de resgate' + notifica pai\n\n"
                 "Alerta Vermelho (0-3 dias)\n"
                 "• Zero sessoes 7+ dias, notificacoes off\n"
                 "→ Email com progresso + oferta premium\n\n"
                 "Resultado: reducao de churn de 5% para 2-3%\n"
                 "ao longo de 24 meses."
             ), accent_color=RED, title_color=RED, body_size=13)

    # Content creation
    add_card(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.5),
             "Criacao Automatizada de Conteudo", (
                 "6 agentes AI em pipeline:\n"
                 "Detector de Gaps → Gerador → Pedagogo Digital\n"
                 "→ Review Humano → A/B Tester → Otimizador\n\n"
                 "AI gera 80% do conteudo, humano valida."
             ), accent_color=TEAL, title_color=TEAL, body_size=14)

    # Stats
    add_stat_card(slide, Inches(6.8), Inches(4.3), "90%", "Reducao no custo\npor licao", GREEN)
    add_stat_card(slide, Inches(9.6), Inches(4.3), "50+", "Licoes novas\npor mes", BLUE)
    add_stat_card(slide, Inches(6.8), Inches(6.0), "5-8x", "Variantes\npor licao", ORANGE)
    add_stat_card(slide, Inches(9.6), Inches(6.0), "3-5", "Pessoas no\ntime (nao 20)", PURPLE)


def build_slide_11_four_buckets_gamification(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "4 Potes + Gamificacao: Educacao com Resultado Real")
    slide_number_badge(slide, "11")

    # 4 Buckets
    buckets = [
        ("Gastar", "Mesada livre,\nconsumo consciente", GREEN),
        ("Guardar", "Objetivos,\npaciencia, metas", BLUE),
        ("Doar", "Generosidade,\nimpacto social", PINK),
        ("Investir", "Futuro, juros,\ncrescimento", ORANGE),
    ]
    for i, (name, desc, color) in enumerate(buckets):
        left = Inches(0.5 + i * 3.15)
        add_rounded_rect(slide, left, Inches(1.8), Inches(2.95), Inches(2.2), CARD_BG, border_color=color, border_width=Pt(1))
        add_textbox(slide, left, Inches(1.95), Inches(2.95), Inches(0.5),
                    name, font_size=20, font_color=color, bold=True,
                    font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.1), Inches(2.5), Inches(2.75), Inches(1.2),
                    desc, font_size=13, font_color=WHITE,
                    font_name=BODY_FONT, alignment=PP_ALIGN.CENTER)

    # Gamification
    add_textbox(slide, Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.4),
                "Gamificacao alimentada por AI — cada desafio calibrado para o perfil individual",
                font_size=16, font_color=WHITE, bold=True, font_name=TITLE_FONT)

    game_items = [
        ("XP Points", "Experiencia por completar\ntarefas e licoes"),
        ("Niveis", "Semente → Broto →\nArvore → Floresta"),
        ("Badges", "Conquistas por marcos\nreais de aprendizado"),
        ("Avatar", "Personalizavel, evolui\ncom o progresso"),
        ("Ranking", "Entre irmaos,\ncompetição saudavel"),
        ("Streaks", "Sequencias que\nmotivam habito"),
    ]
    for i, (title, desc) in enumerate(game_items):
        left = Inches(0.4 + i * 2.1)
        add_rounded_rect(slide, left, Inches(4.9), Inches(1.95), Inches(2.2), CARD_BG_ALT)
        add_textbox(slide, left, Inches(5.0), Inches(1.95), Inches(0.4),
                    title, font_size=13, font_color=GREEN, bold=True,
                    font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.05), Inches(5.5), Inches(1.85), Inches(1.2),
                    desc, font_size=11, font_color=LIGHT_GRAY,
                    font_name=BODY_FONT, alignment=PP_ALIGN.CENTER)


def build_slide_12_whitelabel(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "White-Label: 100% Sua Marca")
    slide_subtitle(slide, "Seus clientes veem sua marca. GRANIX e invisivel — apenas 'by GRANIX' como selo.")
    slide_number_badge(slide, "12")

    # What's customizable
    add_card(slide, Inches(0.5), Inches(2.0), Inches(5.8), Inches(4.8),
             "O que voce customiza", (
                 "• Cores e paleta da sua marca\n"
                 "• Logo e icones\n"
                 "• Fontes tipograficas\n"
                 "• Nome do programa (ex: 'XP Junior')\n"
                 "• Produtos financeiros nos exemplos\n"
                 "  (ex: Tesouro Educa+ da XP, CDB do Inter)\n"
                 "• Tom de comunicacao\n"
                 "• Controle parental branding\n\n"
                 "Integracao: npm install @granix/components\n"
                 "Tempo de setup: 3-4 semanas"
             ), accent_color=GREEN, body_size=14)

    # Integration flow
    add_card(slide, Inches(6.8), Inches(2.0), Inches(6.0), Inches(4.8),
             "Integracao Tecnica", (
                 "Semana 1: Setup\n"
                 "  npm install + configuracao de tema\n"
                 "  Callbacks + documentacao\n\n"
                 "Semana 2-3: Integracao\n"
                 "  Embed de componentes no app\n"
                 "  Conexao de callbacks com APIs internas\n"
                 "  Testes em staging\n\n"
                 "Semana 4: Go-live\n"
                 "  Feature flag para rollout gradual\n"
                 "  Monitoramento conjunto\n\n"
                 "Recurso necessario: 1-2 devs do seu time"
             ), accent_color=BLUE, title_color=BLUE, body_size=13)


def build_slide_13_business_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "Modelo de Negocio: Assinatura + Revenue Share")
    slide_number_badge(slide, "13")

    # Pricing
    plans = [
        ("Gratis", "R$0", "Academy limitada\n1 licao/semana\n2 potes, 3 tarefas/mes", MED_GRAY),
        ("Premium", "R$19,90/mes", "Academy completa com AI\n4 potes, tarefas ilimitadas\nGamificacao basica", GREEN),
        ("Premium+", "R$39,90/mes", "Tudo + AI avancada\nMulti-filhos, simuladores\nRelatorios, desafios exclusivos", ORANGE),
    ]
    for i, (name, price, features, color) in enumerate(plans):
        left = Inches(0.5 + i * 4.1)
        add_rounded_rect(slide, left, Inches(1.5), Inches(3.8), Inches(3.2), CARD_BG, border_color=color, border_width=Pt(2))
        add_textbox(slide, left, Inches(1.65), Inches(3.8), Inches(0.4),
                    name, font_size=18, font_color=color, bold=True,
                    font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, left, Inches(2.1), Inches(3.8), Inches(0.5),
                    price, font_size=28, font_color=WHITE, bold=True,
                    font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.2), Inches(2.7), Inches(3.4), Inches(1.8),
                    features, font_size=13, font_color=LIGHT_GRAY,
                    font_name=BODY_FONT, alignment=PP_ALIGN.CENTER)

    # Revenue share
    add_card(slide, Inches(0.5), Inches(5.0), Inches(5.8), Inches(2.2),
             "Revenue Share", (
                 "Assinatura:  GRANIX 70%  |  Sua instituicao 30%\n"
                 "Transacoes financeiras: 100% para voce\n"
                 "Interchange, float, spread: 100% para voce"
             ), accent_color=GREEN, body_size=14)

    add_card(slide, Inches(6.8), Inches(5.0), Inches(6.0), Inches(2.2),
             "O que voce ganha", (
                 "• 30% de receita recorrente sem custo de build\n"
                 "• CAC intergeracional (jovem 10 → adulto 18)\n"
                 "• Defesa competitiva vs Nubank, Inter, C6\n"
                 "• AI anti-churn reduz abandono de conta familiar"
             ), accent_color=BLUE, title_color=BLUE, body_size=14)


def build_slide_14_compliance(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "Compliance: O Parceiro Mais Seguro Possivel")
    slide_subtitle(slide, "LGPD Art. 14 + ECA Digital (em vigor desde 17/03/2026) + Privacy by Design")
    slide_number_badge(slide, "14")

    # What GRANIX never sees
    add_card(slide, Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.5),
             "GRANIX NUNCA ve", (
                 "❌  Nome, CPF, RG, endereco\n"
                 "❌  Saldo, transacoes, investimentos\n"
                 "❌  Dados de cartao\n"
                 "❌  Biometria ou localizacao"
             ), accent_color=RED, title_color=RED, body_size=15)

    # What GRANIX processes
    add_card(slide, Inches(6.8), Inches(2.0), Inches(6.0), Inches(2.5),
             "GRANIX processa (anonimizado)", (
                 "✅  Faixa etaria (nao idade exata)\n"
                 "✅  Interesses declarados\n"
                 "✅  Comportamento de aprendizado\n"
                 "✅  Progresso de gamificacao"
             ), accent_color=GREEN, body_size=15)

    # ECA Digital compliance
    compliance_items = [
        ("Verificacao de idade", "Seu KYC ja faz"),
        ("Controle parental", "Dashboard do pai nativo"),
        ("Privacy by Design", "Separacao desde dia 1"),
        ("Zero publicidade", "Educacao, nao ads"),
        ("Design seguro", "Zero PII no cloud"),
        ("ISO 27001", "Certificacao Ano 1"),
    ]
    for i, (item, detail) in enumerate(compliance_items):
        col = i % 3
        row = i // 3
        left = Inches(0.5 + col * 4.2)
        top = Inches(4.8 + row * 1.2)
        add_rounded_rect(slide, left, top, Inches(3.9), Inches(1.0), CARD_BG_ALT)
        add_textbox(slide, left + Inches(0.15), top + Inches(0.1), Inches(3.6), Inches(0.4),
                    f"✅  {item}", font_size=14, font_color=GREEN, bold=True, font_name=TITLE_FONT)
        add_textbox(slide, left + Inches(0.15), top + Inches(0.55), Inches(3.6), Inches(0.3),
                    detail, font_size=12, font_color=LIGHT_GRAY, font_name=BODY_FONT)


def build_slide_15_next_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    slide_title(slide, "Proximos Passos: PoC Sem Custo para Voce")
    slide_number_badge(slide, "15")

    # Timeline
    steps = [
        ("AGORA", "Demo personalizada\ncom sua marca", "Custo: R$0", GREEN),
        ("SEMANAS 2-6", "PoC tecnica\nno seu staging", "Custo: R$0", BLUE),
        ("MESES 2-4", "Piloto com\n500-1.000 familias", "Metricas reais", ORANGE),
        ("MES 4", "Go / No-Go\nbaseado em dados", "Decisao informada", PURPLE),
        ("MES 5+", "Rollout\npara base completa", "Receita recorrente", GREEN),
    ]
    for i, (phase, desc, detail, color) in enumerate(steps):
        left = Inches(0.3 + i * 2.55)
        top = Inches(1.8)
        w = Inches(2.35)
        add_rounded_rect(slide, left, top, w, Inches(3.5), CARD_BG, border_color=color, border_width=Pt(1))
        add_textbox(slide, left, top + Inches(0.2), w, Inches(0.4),
                    phase, font_size=16, font_color=color, bold=True,
                    font_name=TITLE_FONT, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, left, top + Inches(0.8), w, Inches(1.2),
                    desc, font_size=14, font_color=WHITE,
                    font_name=BODY_FONT, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, left, top + Inches(2.3), w, Inches(0.5),
                    detail, font_size=12, font_color=MED_GRAY,
                    font_name=BODY_FONT, alignment=PP_ALIGN.CENTER)

        # Arrow between steps
        if i < len(steps) - 1:
            add_textbox(slide, left + w - Inches(0.1), top + Inches(1.3), Inches(0.5), Inches(0.4),
                        "→", font_size=24, font_color=MED_GRAY,
                        font_name=BODY_FONT, alignment=PP_ALIGN.CENTER)

    # Key message
    add_multiline_textbox(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(1.5), [
        ("Risco zero para sua instituicao:", 20, GREEN, True, TITLE_FONT),
        ("PoC e piloto sem custo. Decisao baseada em dados reais.", 18, WHITE, False, BODY_FONT),
        ("Voce so paga quando funciona.", 18, WHITE, False, BODY_FONT),
        ("", 10, DARK_BG, False, BODY_FONT),
        ("contato@granix.com.br  |  granix.com.br", 14, MED_GRAY, False, BODY_FONT),
    ], alignment=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_01_cover(prs)
    build_slide_02_gap(prs)
    build_slide_03_barrier(prs)
    build_slide_04_solution(prs)
    build_slide_05_onboarding(prs)
    build_slide_06_first_session(prs)
    build_slide_07_weekly_loop(prs)
    build_slide_08_content(prs)
    build_slide_09_ai_personalization(prs)
    build_slide_10_ai_prediction(prs)
    build_slide_11_four_buckets_gamification(prs)
    build_slide_12_whitelabel(prs)
    build_slide_13_business_model(prs)
    build_slide_14_compliance(prs)
    build_slide_15_next_steps(prs)

    prs.save(OUTPUT_PATH)
    print(f"Pitch deck saved to {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
